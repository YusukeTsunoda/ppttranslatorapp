from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from typing import Dict, List
import json
import sys

def apply_text_style(run, style: Dict):
    """テキストスタイルを適用"""
    if "fontSize" in style:
        run.font.size = Pt(style["fontSize"])
    if "fontFamily" in style:
        run.font.name = style["fontFamily"]
    if "bold" in style:
        run.font.bold = style["bold"]
    if "italic" in style:
        run.font.italic = style["italic"]
    if "color" in style:
        # RGB形式の色を設定
        color = style["color"]
        run.font.color.rgb = color

def pixels_to_emu(pixels: float) -> int:
    """ピクセルをEMU単位に変換"""
    EMU_PER_INCH = 914400
    PIXELS_PER_INCH = 96
    return int(pixels * EMU_PER_INCH / PIXELS_PER_INCH)

def create_translated_pptx(original_pptx_path: str, translations: List[Dict], output_path: str):
    """翻訳済みのPPTXファイルを生成"""
    try:
        # 元のプレゼンテーションを読み込み
        prs = Presentation(original_pptx_path)
        
        # 各スライドの翻訳を適用
        for slide_idx, slide_data in enumerate(translations):
            if slide_idx >= len(prs.slides):
                print(f"Warning: Slide {slide_idx + 1} not found in original presentation")
                continue
                
            slide = prs.slides[slide_idx]
            texts = slide_data.get("texts", [])
            
            # 既存のテキストボックスを更新
            shape_index = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape_index < len(texts):
                    text_data = texts[shape_index]
                    
                    # テキストボックスの位置とサイズを更新
                    position = text_data.get("position", {})
                    shape.left = pixels_to_emu(position.get("x", 0) * 2)  # スケール係数を元に戻す
                    shape.top = pixels_to_emu(position.get("y", 0) * 2)   # スケール係数を元に戻す
                    shape.width = pixels_to_emu(position.get("width", 0) * 2)
                    shape.height = pixels_to_emu(position.get("height", 0) * 2)
                    
                    # テキストを更新
                    if shape.text_frame:
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = text_data.get("translation", "") or text_data.get("text", "")
                        
                        # スタイルを適用
                        if "style" in text_data:
                            apply_text_style(run, text_data["style"])
                    
                    shape_index += 1
        
        # 変更を保存
        prs.save(output_path)
        print(json.dumps({
            "success": True,
            "message": "Successfully generated translated PPTX",
            "output_path": output_path
        }))
        sys.exit(0)
        
    except Exception as e:
        error_result = {
            "success": False,
            "error": str(e)
        }
        print(json.dumps(error_result))
        sys.exit(1)

if __name__ == "__main__":
    # コマンドライン引数からパラメータを取得
    if len(sys.argv) != 4:
        print(json.dumps({
            "success": False,
            "error": "Invalid arguments. Required: original_pptx_path, translations_json_path, output_path"
        }))
        sys.exit(1)
    
    original_pptx_path = sys.argv[1]
    translations_json_path = sys.argv[2]
    output_path = sys.argv[3]
    
    # 翻訳データを読み込み
    with open(translations_json_path, 'r', encoding='utf-8') as f:
        translations = json.load(f)
    
    create_translated_pptx(original_pptx_path, translations, output_path)
