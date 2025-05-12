from pptx import Presentation
import json
import sys
import uuid
from typing import Dict, List, Optional

# EMUをピクセルに変換する定数
EMU_PER_INCH = 914400  # 1インチ = 914400 EMU
PIXELS_PER_INCH = 96   # 1インチ = 96ピクセル

# 16:9のスライドサイズ
WIDESCREEN_WIDTH = 720  # 16:9の幅
WIDESCREEN_HEIGHT = 405  # 16:9の高さ

# 4:3のスライドサイズ
STANDARD_WIDTH = 720  # 4:3の幅
STANDARD_HEIGHT = 540  # 4:3の高さ

def emu_to_pixels(emu: int) -> float:
    """EMU単位をピクセルに変換"""
    return (emu * PIXELS_PER_INCH) / EMU_PER_INCH

def get_slide_dimensions(presentation_width: int, presentation_height: int) -> Dict[str, int]:
    """プレゼンテーションのアスペクト比に基づいて適切なサイズを返す"""
    aspect_ratio = presentation_width / presentation_height
    print(f"Debug - Presentation size: {emu_to_pixels(presentation_width)}x{emu_to_pixels(presentation_height)} pixels")
    print(f"Debug - Aspect ratio: {aspect_ratio}")
    
    # 16:9に近いアスペクト比の場合
    if abs(aspect_ratio - 16/9) < abs(aspect_ratio - 4/3):
        return {
            "width": WIDESCREEN_WIDTH,
            "height": WIDESCREEN_HEIGHT
        }
    # 4:3に近いアスペクト比の場合
    else:
        return {
            "width": STANDARD_WIDTH,
            "height": STANDARD_HEIGHT
        }

def normalize_position(position: Dict, presentation_width: int, presentation_height: int) -> Dict:
    """位置情報を標準サイズに正規化"""
    # プレゼンテーションのアスペクト比に基づいて適切なサイズを取得
    target_dimensions = get_slide_dimensions(presentation_width, presentation_height)
    
    # 実際のサイズをピクセルに変換
    actual_width = emu_to_pixels(presentation_width)
    actual_height = emu_to_pixels(presentation_height)
    
    # スケール係数を計算
    width_ratio = target_dimensions["width"] / actual_width
    height_ratio = target_dimensions["height"] / actual_height
    
    # 位置情報を変換
    normalized = {
        "x": round(emu_to_pixels(position["x"]) * width_ratio * 0.5),  # スケール係数0.5を適用
        "y": round(emu_to_pixels(position["y"]) * height_ratio * 1.66),  # スケール係数0.5を適用
        "width": round(emu_to_pixels(position["width"]) * width_ratio * 0.5),  # スケール係数0.5を適用
        "height": round(emu_to_pixels(position["height"]) * height_ratio * 0.5)  # スケール係数0.5を適用
    }
    
    print(f"Debug - Original position (EMU): {position}")
    print(f"Debug - Converted position (px): {normalized}")
    
    return normalized

def extract_text_style(run) -> Dict:
    """テキストのスタイル情報を抽出"""
    return {
        "fontSize": run.font.size.pt if hasattr(run.font, "size") and run.font.size else None,
        "fontFamily": run.font.name if hasattr(run.font, "name") else None,
        "bold": run.font.bold if hasattr(run.font, "bold") else None,
        "italic": run.font.italic if hasattr(run.font, "italic") else None,
    }

def extract_position(shape, presentation_width: int, presentation_height: int) -> Dict:
    """シェイプの位置情報を抽出"""
    raw_position = {
        "x": shape.left,
        "y": shape.top,
        "width": shape.width,
        "height": shape.height
    }
    return normalize_position(raw_position, presentation_width, presentation_height)

def extract_slide_content(slide, presentation_width: int, presentation_height: int) -> Dict:
    """スライドの内容を抽出"""
    texts = []
    layout_elements = []
    
    # スライドの実際のサイズをピクセルに変換
    slide_width = emu_to_pixels(presentation_width)
    slide_height = emu_to_pixels(presentation_height)
    
    print(f"Debug - Processing slide - Original size: {slide_width}x{slide_height} pixels")

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # EMUからピクセルに変換し、さらにスケーリング
            raw_position = {
                "x": round(emu_to_pixels(shape.left) * 0.5),  # スケール係数0.5を適用
                "y": round(emu_to_pixels(shape.top) * 0.5),   # スケール係数0.5を適用
                "width": round(emu_to_pixels(shape.width) * 0.5),  # スケール係数0.5を適用
                "height": round(emu_to_pixels(shape.height) * 0.5)  # スケール係数0.5を適用
            }
            
            print(f"Debug - Text shape: '{shape.text.strip()}'")
            print(f"Debug - Original position (EMU): left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
            print(f"Debug - Scaled position (px): {raw_position}")
            
            # テキスト要素の抽出
            text_element = {
                "id": str(uuid.uuid4()),
                "text": shape.text.strip(),
                "position": raw_position,  # ← x, y, width, height を必ず含める
                "style": extract_text_style(shape.text_frame.paragraphs[0].runs[0]) if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs else {}
            }
            texts.append(text_element)

        # レイアウト要素の抽出
        layout_element = {
            "type": shape.shape_type,
            "position": extract_position(shape, presentation_width, presentation_height),
        }
        if hasattr(shape, "placeholder_format"):
            layout_element["placeholder"] = shape.placeholder_format.type
        layout_elements.append(layout_element)

    return {
        "texts": texts,
        "layout": {
            "masterLayout": slide.slide_layout.name,
            "elements": layout_elements
        }
    }

def extract_metadata(presentation) -> Dict:
    """プレゼンテーションのメタデータを抽出"""
    core_properties = presentation.core_properties
    
    # オリジナルサイズをピクセルに変換
    original_width = emu_to_pixels(presentation.slide_width)
    original_height = emu_to_pixels(presentation.slide_height)
    
    # アスペクト比を計算
    aspect_ratio = original_width / original_height
    
    print(f"Debug - Original size in EMU: {presentation.slide_width}x{presentation.slide_height}")
    print(f"Debug - Original size in pixels: {original_width}x{original_height}")
    print(f"Debug - Aspect ratio: {aspect_ratio}")
    
    # 基準サイズを設定（16:9 = 720x405, 4:3 = 720x540）
    if abs(aspect_ratio - 16/9) < abs(aspect_ratio - 4/3):
        target_width = 720
        target_height = 405
    else:
        target_width = 720
        target_height = 540
    
    metadata = {
        "totalSlides": len(presentation.slides),
        "title": core_properties.title if hasattr(core_properties, "title") else None,
        "author": core_properties.author if hasattr(core_properties, "author") else None,
        "lastModified": core_properties.modified.isoformat() if hasattr(core_properties, "modified") and core_properties.modified else None,
        "dimensions": {
            "width": target_width,
            "height": target_height
        },
        "aspectRatio": "16:9" if abs(aspect_ratio - 16/9) < 0.1 else "4:3",
        "originalSize": {
            "width": round(original_width),
            "height": round(original_height)
        }
    }
    
    print(f"Debug - Final metadata: {metadata}")
    return metadata

def main(pptx_path: str):
    """メイン処理"""
    try:
        prs = Presentation(pptx_path)
        slides = []
        
        # プレゼンテーションのサイズを取得
        presentation_width = prs.slide_width
        presentation_height = prs.slide_height
        
        # メタデータを先に生成
        metadata = extract_metadata(prs)
        print(f"Debug - Generated metadata: {metadata}")
        
        for idx, slide in enumerate(prs.slides):
            content = extract_slide_content(slide, presentation_width, presentation_height)
            slide_data = {
                "id": str(uuid.uuid4()),
                "index": idx,
                "metadata": metadata,  # 各スライドにメタデータを追加
                **content
            }
            slides.append(slide_data)
        
        result = {
            "success": True,
            "data": {
                "slides": slides,
                "metadata": metadata  # ルートレベルにもメタデータを追加
            }
        }
        
        print(json.dumps(result))
        sys.exit(0)
    
    except Exception as e:
        error_result = {
            "success": False,
            "error": str(e)
        }
        print(json.dumps(error_result))
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({
            "success": False,
            "error": "PPTXファイルのパスを指定してください"
        }))
        sys.exit(1)
    
    main(sys.argv[1])
