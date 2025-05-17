from pathlib import Path
import random
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
import json

class TestDataGenerator:
    def __init__(self, output_dir: str = "tests/performance/test_data"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self._load_sample_texts()

    def _load_sample_texts(self):
        """サンプルテキストの読み込み"""
        sample_texts_path = Path(__file__).parent / "sample_texts.json"
        try:
            with open(sample_texts_path, "r", encoding="utf-8") as f:
                self.sample_texts = json.load(f)
        except FileNotFoundError:
            print(f"Warning: {sample_texts_path} not found. Using default sample texts.")
            self.sample_texts = {
                "ja": ["これはテストテキストです。"],
                "en": ["This is a test text."]
            }

    def generate_test_pptx(
        self,
        slide_count: int,
        text_boxes_per_slide: range,
        languages: List[str] = ["ja", "en"],
        filename: str = None,
        font_sizes: range = range(10, 25),
        text_box_styles: List[Dict] = None
    ) -> Path:
        """テスト用PPTXファイルを生成"""
        prs = Presentation()
        metadata = {
            "slide_count": slide_count,
            "text_boxes": [],
            "languages": languages,
            "font_size_range": {"min": font_sizes.start, "max": font_sizes.stop}
        }

        # デフォルトのテキストボックススタイル
        default_styles = [
            {"width": (2.0, 4.0), "height": (0.5, 2.0)},  # 標準サイズ
            {"width": (4.0, 6.0), "height": (2.0, 3.0)},  # 大きいサイズ
            {"width": (1.0, 2.0), "height": (0.3, 1.0)}   # 小さいサイズ
        ]
        text_box_styles = text_box_styles or default_styles

        for i in range(slide_count):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
            text_box_count = random.randint(text_boxes_per_slide.start, text_boxes_per_slide.stop)
            
            for j in range(text_box_count):
                # スタイルをランダムに選択
                style = random.choice(text_box_styles)
                width_range = style["width"]
                height_range = style["height"]
                
                # ランダムな位置とサイズでテキストボックスを配置
                left = random.uniform(0.5, 8.5)
                top = random.uniform(0.5, 6.5)
                width = random.uniform(*width_range)
                height = random.uniform(*height_range)
                
                textbox = slide.shapes.add_textbox(
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
                
                # ランダムな言語とテキストを選択
                lang = random.choice(languages)
                text = random.choice(self.sample_texts[lang])
                
                # テキストボックスの設定
                text_frame = textbox.text_frame
                text_frame.text = text
                text_frame.paragraphs[0].font.size = Pt(random.randint(font_sizes.start, font_sizes.stop))
                
                # メタデータを記録
                metadata["text_boxes"].append({
                    "slide": i,
                    "position": {"left": left, "top": top},
                    "size": {"width": width, "height": height},
                    "language": lang,
                    "text_length": len(text),
                    "font_size": text_frame.paragraphs[0].font.size.pt
                })

        # ファイル名が指定されていない場合は自動生成
        if filename is None:
            filename = f"test_slides_{slide_count}_{len(languages)}.pptx"
        
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        
        # メタデータを保存
        metadata_path = output_path.with_suffix(".json")
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)
        
        return output_path

    def generate_edge_cases(self) -> List[Path]:
        """エッジケースのテストファイルを生成"""
        edge_cases = []
        
        # 1. 極端に多いスライド数
        edge_cases.append(self.generate_test_pptx(
            100, range(1, 3), filename="edge_many_slides.pptx"
        ))
        
        # 2. スライドあたり多数のテキストボックス
        edge_cases.append(self.generate_test_pptx(
            5, range(20, 30), filename="edge_many_textboxes.pptx"
        ))
        
        # 3. 極端なフォントサイズ
        edge_cases.append(self.generate_test_pptx(
            3, range(2, 4),
            font_sizes=range(36, 72),
            filename="edge_large_fonts.pptx"
        ))
        
        # 4. 極端なテキストボックスサイズ
        large_styles = [
            {"width": (6.0, 8.0), "height": (4.0, 5.0)},  # 非常に大きい
            {"width": (0.5, 1.0), "height": (0.2, 0.5)}   # 非常に小さい
        ]
        edge_cases.append(self.generate_test_pptx(
            5, range(2, 5),
            text_box_styles=large_styles,
            filename="edge_textbox_sizes.pptx"
        ))
        
        # 5. 多言語混在
        edge_cases.append(self.generate_test_pptx(
            5, range(2, 5),
            languages=["ja", "en", "zh"],
            filename="edge_multilingual.pptx"
        ))
        
        return edge_cases

    def generate_dataset(
        self,
        dataset_name: str,
        configurations: List[Dict]
    ) -> Dict[str, List[Path]]:
        """テストデータセットを生成"""
        dataset_dir = self.output_dir / dataset_name
        dataset_dir.mkdir(exist_ok=True)
        
        generated_files = {}
        for config in configurations:
            category = config.pop("category", "default")
            if category not in generated_files:
                generated_files[category] = []
            
            # 設定に基づいてファイルを生成
            filename = f"{category}_{len(generated_files[category])}.pptx"
            file_path = self.generate_test_pptx(
                filename=filename,
                **config
            )
            generated_files[category].append(file_path)
        
        # データセット設定を保存
        config_path = dataset_dir / "dataset_config.json"
        with open(config_path, "w", encoding="utf-8") as f:
            json.dump({
                "name": dataset_name,
                "configurations": configurations,
                "files": {k: [str(p) for p in v] for k, v in generated_files.items()}
            }, f, indent=2)
        
        return generated_files 