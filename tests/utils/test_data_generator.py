import random
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import lorem
import string

def generate_random_text(length: int) -> str:
    """ランダムなテキストを生成"""
    return ''.join(random.choices(string.ascii_letters + string.digits + ' ', k=length))

def generate_test_pptx(
    slide_count: int = 10,
    text_density: str = 'normal',
    image_count: int = 0,
    corrupt: bool = False
) -> str:
    """テスト用のPPTXファイルを生成"""
    prs = Presentation()
    
    # テキスト密度に応じてテキスト長を決定
    text_lengths = {
        'empty': 0,
        'low': 50,
        'normal': 200,
        'high': 1000
    }
    text_length = text_lengths.get(text_density, 200)
    
    # スライドの生成
    for _ in range(slide_count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # タイトルとコンテンツのレイアウト
        
        # タイトルの追加
        if slide.shapes.title:
            slide.shapes.title.text = generate_random_text(20)
        
        # 本文の追加
        if text_length > 0:
            body_shape = slide.shapes.placeholders[1]
            body_shape.text = generate_random_text(text_length)
    
    # ファイルの保存
    output_dir = Path('tests/performance/test_data')
    output_dir.mkdir(parents=True, exist_ok=True)
    
    file_path = output_dir / f'test_{slide_count}_{text_density}.pptx'
    prs.save(str(file_path))
    
    # 破損ファイルの生成（テスト用）
    if corrupt:
        with open(file_path, 'r+b') as f:
            f.seek(random.randint(0, 100))
            f.write(b'corrupt')
    
    return str(file_path)

def generate_multilingual_test_pptx(slide_count: int = 10, output_dir: str = "tests/test_data") -> str:
    """
    多言語テスト用のPPTXファイルを生成
    
    Args:
        slide_count: 生成するスライドの数
        output_dir: 出力ディレクトリ
        
    Returns:
        生成されたファイルのパス
    """
    # 日本語のサンプルテキスト
    japanese_texts = [
        "こんにちは、世界",
        "プレゼンテーションのテスト",
        "自動化されたテストケース",
        "パフォーマンス最適化",
        "多言語対応"
    ]
    
    # 出力ディレクトリの作成
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドの生成
    for i in range(slide_count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # タイトルの設定（日本語）
        title = slide.shapes.title
        title.text = japanese_texts[i % len(japanese_texts)]
        
        # 英語のコンテンツ
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = lorem.paragraph()
        
        # 日本語のテキストボックス
        left = Inches(random.uniform(1, 8))
        top = Inches(random.uniform(1, 5))
        txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(1))
        tf = txBox.text_frame
        tf.text = japanese_texts[(i + 1) % len(japanese_texts)]
    
    # ファイルの保存
    output_path = f"{output_dir}/multilingual_test_{slide_count}.pptx"
    prs.save(output_path)
    
    return output_path 