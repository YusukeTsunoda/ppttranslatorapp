#!/usr/bin/env python3
import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

def find_matching_shape(shape, original_text):
    """シェイプのテキストが元のテキストと一致するか確認"""
    if not hasattr(shape, 'text_frame'):
        return False
    return shape.text_frame.text.strip() == original_text.strip()

def adjust_text_frame_for_language(text_frame, translation, language="en"):
    """言語に応じてテキストフレームの設定を調整"""
    # 言語別の調整係数
    lang_factors = {
        "en": 1.0,    # 英語（基準）
        "ja": 1.2,    # 日本語
        "zh": 1.2,    # 中国語
        "ko": 1.15,   # 韓国語
        "ru": 1.1,    # ロシア語
        "ar": 1.1,    # アラビア語
        "de": 1.1,    # ドイツ語
        "fr": 1.05,   # フランス語
        "es": 1.05,   # スペイン語
        "it": 1.05,   # イタリア語
        "pt": 1.05,   # ポルトガル語
    }
    
    # 言語係数（デフォルト: 1.0）
    lang_factor = lang_factors.get(language, 1.0)
    
    try:
        # テキストの自動サイズ調整
        if len(translation) > len(text_frame.text) * 1.2:
            # 翻訳テキストが元のテキストより20%以上長い場合は自動サイズ調整を有効に
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        else:
            # それ以外の場合はテキストを枠内に収める
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # 縦方向のテキスト配置を中央に
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # テキストの折り返しを有効に
        text_frame.word_wrap = True
        
        # 余白の調整（特に日本語などの場合）
        if lang_factor > 1.0:
            # 余白を少し縮小
            margin = Emu(46752)  # 約0.064インチ（通常の2/3程度）
            text_frame.margin_left = margin
            text_frame.margin_right = margin
            text_frame.margin_top = margin
            text_frame.margin_bottom = margin
    except Exception as e:
        print(f"  - Warning: テキストフレーム調整エラー: {e}", file=sys.stderr)

def detect_language(text):
    """簡易的な言語検出 - 最初の数文字から言語を推定"""
    # 日本語、中国語、韓国語の文字コード範囲
    if any(ord(c) > 0x3000 for c in text[:20]):
        # 日本語の特徴的な文字が含まれているか
        if any(0x3040 <= ord(c) <= 0x309F or 0x30A0 <= ord(c) <= 0x30FF for c in text[:20]):
            return "ja"
        # 韓国語の特徴的な文字が含まれているか
        elif any(0xAC00 <= ord(c) <= 0xD7A3 for c in text[:20]):
            return "ko"
        # それ以外はとりあえず中国語と判定
        else:
            return "zh"
    
    # キリル文字（ロシア語など）
    elif any(0x0400 <= ord(c) <= 0x04FF for c in text[:20]):
        return "ru"
    
    # アラビア文字
    elif any(0x0600 <= ord(c) <= 0x06FF for c in text[:20]):
        return "ar"
    
    # 西欧諸国の言語は文字だけでは判別が難しいので、一般的な特徴で判定
    # ドイツ語の特徴的な文字
    elif any(c in "äöüß" for c in text[:50]):
        return "de"
    
    # フランス語の特徴的な文字
    elif any(c in "éèêëàâçùûüÿæœ" for c in text[:50]):
        return "fr"
    
    # スペイン語の特徴的な文字
    elif any(c in "áéíóúüñ¿¡" for c in text[:50]):
        return "es"
    
    # イタリア語の特徴的な文字
    elif any(c in "àèéìíîòóùú" for c in text[:50]) and "che " in text[:50].lower():
        return "it"
    
    # ポルトガル語の特徴的な文字
    elif any(c in "áàâãéêíóôõúç" for c in text[:50]):
        return "pt"
    
    # デフォルトは英語
    return "en"

def adjust_font_size(paragraph, min_size=8, max_size=None):
    """テキストのフォントサイズを二分探索で最適化する"""
    
    # 現在のフォントサイズを取得
    current_size = None
    for run in paragraph.runs:
        if run.font.size is not None:
            current_size = run.font.size.pt
            break
    
    # 現在のサイズが取得できなかった場合はデフォルト値を使用
    if current_size is None:
        current_size = 12
    
    # 最大サイズが指定されていない場合は現在のサイズを使用
    if max_size is None:
        max_size = current_size
    
    def check_text_fits(size):
        """指定されたフォントサイズでテキストが収まるかチェック"""
        # テキストフレームの幅と高さを取得
        frame_width = paragraph._parent.width
        frame_height = paragraph._parent.height
        
        # テキストの予想サイズを計算
        text = paragraph.text
        estimated_width = len(text) * size * 0.6  # 文字あたりの平均幅を概算
        estimated_height = size * 1.2  # 行の高さを概算
        
        # 行数を計算
        lines = estimated_width / frame_width
        total_height = lines * estimated_height
        
        return total_height <= frame_height
    
    # 二分探索でフォントサイズを決定
    left = min_size
    right = max_size
    optimal_size = min_size
    
    while left <= right:
        mid = (left + right) / 2
        if check_text_fits(mid):
            optimal_size = mid
            left = mid + 0.5  # より大きいサイズを試す
        else:
            right = mid - 0.5  # より小さいサイズを試す
        
        # 十分な精度に達したら終了
        if right - left < 0.5:
            break
    
    # 最適なサイズを設定
    for run in paragraph.runs:
        run.font.size = Pt(optimal_size)

def update_slide_text(slide, slide_data):
    """スライド内のテキストを翻訳文で更新する"""
    if not slide_data.get('texts'):
        return

    for i, text_data in enumerate(slide_data['texts']):
        original_text = text_data.get('text', '').strip()
        
        # 翻訳テキストの取得方法を修正
        translation = None
        # ターゲット言語を取得
        target_language = slide_data.get('targetLanguage', 'en')
        
        # 1. translations配列から取得を試みる
        if slide_data.get('translations') and i < len(slide_data['translations']):
            trans_entry = slide_data['translations'][i]
            if isinstance(trans_entry, dict) and trans_entry.get('text'):
                translation = trans_entry.get('text')
            elif isinstance(trans_entry, str):
                translation = trans_entry
        
        # 2. text_dataから直接translationを取得
        if not translation and text_data.get('translation'):
            translation = text_data.get('translation', '')
            
        # 翻訳が存在しない場合はスキップ
        if not translation:
            print(f"No translation found for text: {original_text}", file=sys.stderr)
            continue

        print(f"Updating text: '{original_text}' -> '{translation}'", file=sys.stderr)
        
        # 言語の検出（ターゲット言語が指定されていない場合）
        if not target_language or target_language == 'auto':
            target_language = detect_language(translation)
            print(f"  - Detected language: {target_language}", file=sys.stderr)
        
        # スライド内の全シェイプを検索
        found = False
        for shape in slide.shapes:
            if find_matching_shape(shape, original_text):
                try:
                    # テキストフレームにアクセス
                    if not hasattr(shape, 'text_frame'):
                        print(f"  - Warning: Shape does not have a text frame", file=sys.stderr)
                        continue
                    
                    # テキストフレームを言語に応じて調整
                    adjust_text_frame_for_language(shape.text_frame, translation, target_language)
                    
                    # テキストを更新
                    # まず既存のテキストを保持
                    existing_text = shape.text_frame.text
                    
                    # テキストを設定
                    if shape.text_frame.paragraphs:
                        # 既存の段落のプロパティを保持
                        p = shape.text_frame.paragraphs[0]
                        
                        # フォントの配置情報を保持
                        alignment = p.alignment
                        
                        # 段落のテキストをクリア
                        for run in list(p.runs):
                            p._element.remove(run._element)
                        
                        # 新しいランを追加
                        run = p.add_run()
                        run.text = translation
                        
                        # フォントサイズを調整
                        adjust_font_size(p)
                        
                        # 配置を復元
                        p.alignment = alignment
                    else:
                        # 段落がない場合は単純にテキストを設定
                        shape.text_frame.text = translation
                    
                    print(f"  - Text updated successfully", file=sys.stderr)
                    found = True
                    break
                except Exception as e:
                    print(f"  - Error updating text: {e}", file=sys.stderr)
                
        if not found:
            print(f"  - Warning: No matching shape found for text: {original_text}", file=sys.stderr)

def main():
    if len(sys.argv) != 4:
        print(json.dumps({
            "error": "Usage: python pptx_translator.py input_file output_file slides_data"
        }))
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    slides_data = json.loads(sys.argv[3])

    try:
        # デバッグ情報を出力
        print(f"Loading presentation from: {input_file}", file=sys.stderr)
        
        # 元のプレゼンテーションを読み込む
        prs = Presentation(input_file)
        print(f"Loaded presentation with {len(prs.slides)} slides", file=sys.stderr)

        # スライドデータの構造をデバッグ出力
        print(f"Slides data structure: {json.dumps(slides_data[0] if slides_data else {}, indent=2)}", file=sys.stderr)

        # 各スライドの翻訳を適用
        for slide_data in slides_data:
            slide_index = slide_data.get('index', 0)
            if slide_index < len(prs.slides):
                print(f"\nProcessing slide {slide_index}", file=sys.stderr)
                update_slide_text(prs.slides[slide_index], slide_data)

        # 翻訳済みプレゼンテーションを保存
        print(f"\nSaving presentation to: {output_file}", file=sys.stderr)
        prs.save(output_file)
        print(json.dumps({"success": True}))

    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        print(json.dumps({"error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main() 