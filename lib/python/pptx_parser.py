#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import argparse
from pdf2image import convert_from_path
import subprocess
import logging
from typing import List, Dict, Any, Tuple, Optional, Union
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderGraphicFrame
# 同じディレクトリにあるimage_optimizerモジュールをインポート
from lib.python.image_optimizer import optimize_image, batch_optimize, get_optimal_format

# ロギング設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('pptx_parser')

def convert_to_png(pptx_path: str, output_dir: str, optimize: bool = True, 
                  format: str = 'WEBP', quality: int = 85, 
                  max_width: int = 1920, max_height: int = 1080) -> tuple[List[str], tuple[int, int]]:
    """
    PPTXの各スライドをPNG画像に変換し、必要に応じて最適化する
    
    Args:
        pptx_path (str): PPTXファイルのパス
        output_dir (str): 出力ディレクトリ
        optimize (bool): 画像最適化を行うかどうか
        format (str): 出力画像形式 ('PNG', 'WEBP', 'JPEG')
        quality (int): 画質 (0-100)
        max_width (int): 最大幅
        max_height (int): 最大高さ
    
    Returns:
        tuple[List[str], tuple[int, int]]: 画像パスのリストと画像サイズ
    """
    try:
        # デバッグ出力を追加
        logger.info(f"Starting convert_to_png: pptx_path={pptx_path}, output_dir={output_dir}")
        
        # 出力ディレクトリが存在しない場合は作成
        os.makedirs(output_dir, exist_ok=True)
        
        # PPTXをPDFに変換
        pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, f"{pptx_name}.pdf")
        
        logger.info(f"Converting PPTX to PDF: output={pdf_path}")
        
        if sys.platform == "darwin":  # macOS
            # LibreOfficeのパスを確認
            libreoffice_path = subprocess.run(["which", "soffice"], capture_output=True, text=True).stdout.strip()
            logger.info(f"LibreOffice path: {libreoffice_path}")
            
            result = subprocess.run(
                [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
                capture_output=True,
                text=True
            )
            
            logger.info(f"LibreOffice conversion result: returncode={result.returncode}")
            logger.debug(f"LibreOffice stdout: {result.stdout}")
            logger.debug(f"LibreOffice stderr: {result.stderr}")
            
            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed with code {result.returncode}")
                return [], (0, 0)
        else:  # Linux
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
                capture_output=True,
                text=True
            )
            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed with code {result.returncode}")
                return [], (0, 0)

        if not os.path.exists(pdf_path):
            logger.error(f"PDF file not found at {pdf_path}")
            return [], (0, 0)
            
        logger.info(f"PDF file created successfully: {pdf_path}")
            
        # PDFを画像に変換（アスペクト比を保持）
        try:
            # popplerのパスを確認
            poppler_path = None
            if sys.platform == "darwin":  # macOS
                # Homebrewでインストールされたpopplerのパスを使用
                poppler_path = "/opt/homebrew/bin"
                if not os.path.exists(poppler_path):
                    poppler_path = "/usr/local/bin"
                
                # パスの存在確認
                logger.debug(f"Checking poppler path: {poppler_path}")
                if os.path.exists(os.path.join(poppler_path, "pdftoppm")):
                    logger.debug(f"pdftoppm found at {os.path.join(poppler_path, 'pdftoppm')}")
                else:
                    logger.warning(f"pdftoppm NOT found at {os.path.join(poppler_path, 'pdftoppm')}")
                    # 環境変数PATHからpoppler関連コマンドを探す
                    for path_dir in os.environ.get('PATH', '').split(':'):
                        if os.path.exists(os.path.join(path_dir, 'pdftoppm')):
                            poppler_path = path_dir
                            logger.info(f"Found pdftoppm in PATH: {os.path.join(path_dir, 'pdftoppm')}")
                            break
            
            logger.info(f"Using poppler_path: {poppler_path}")
            logger.info(f"Converting PDF to images with target size: {max_width}x{max_height}")
            
            # 高品質設定でPDFを画像に変換
            images = convert_from_path(
                pdf_path, 
                size=(max_width, max_height),
                poppler_path=poppler_path,
                dpi=300,  # 高解像度
                fmt="png",  # PNG形式で出力（後で最適化）
                transparent=False,  # 透明度なし
                use_cropbox=True,  # クロップボックスを使用
                strict=False  # 厳密なエラーチェックを無効化
            )
            
            logger.info(f"Converted {len(images)} images from PDF")
            
            if not images:
                logger.error("No images were converted from PDF")
                return [], (0, 0)
                
            # 実際の画像サイズを取得
            actual_width = images[0].width
            actual_height = images[0].height
            logger.info(f"Actual image size: {actual_width}x{actual_height}")
            
        except Exception as e:
            logger.error(f"Error converting PDF to images: {str(e)}")
            return [], (0, 0)
        
        # 一時的な画像パスを生成（最適化前）
        temp_image_paths = []
        for i, image in enumerate(images):
            temp_image_path = os.path.join(output_dir, f"temp_slide_{i+1}.png")
            try:
                # 保存前にディレクトリが存在することを確認
                os.makedirs(os.path.dirname(temp_image_path), exist_ok=True)
                
                # 一時的に高品質PNGで保存
                image.save(temp_image_path, "PNG", optimize=True)
                
                if os.path.exists(temp_image_path):
                    logger.debug(f"Saved temporary image: {temp_image_path}")
                    temp_image_paths.append(temp_image_path)
                else:
                    logger.warning(f"Warning: Temporary image file not found: {temp_image_path}")
            except Exception as e:
                logger.error(f"Error saving temporary image {i+1}: {str(e)}")
        
        # 画像の最適化処理
        image_paths = []
        if optimize and temp_image_paths:
            logger.info(f"Optimizing {len(temp_image_paths)} images to format: {format}, quality: {quality}")
            
            # 最適化結果を格納するディレクトリ
            optimized_dir = os.path.join(output_dir, "optimized")
            os.makedirs(optimized_dir, exist_ok=True)
            
            # 画像を最適化
            optimization_results = batch_optimize(
                input_paths=temp_image_paths,
                output_dir=optimized_dir,
                format=format,
                quality=quality,
                max_width=max_width,
                max_height=max_height
            )
            
            # 最適化された画像パスを収集
            for i, result in enumerate(optimization_results):
                if 'error' not in result:
                    # 相対パスを生成
                    rel_path = os.path.relpath(result['path'], output_dir)
                    image_paths.append(rel_path)
                    logger.info(f"Optimized image {i+1}: {result['path']}")
                    logger.info(f"  Size reduction: {result['original']['size']} -> {result['optimized']['size']} bytes")
                    logger.info(f"  Compression ratio: {result['optimized']['compression_ratio']:.2f}%")
                else:
                    logger.error(f"Error optimizing image {i+1}: {result.get('error', 'Unknown error')}")
            
            # 一時ファイルを削除
            for temp_path in temp_image_paths:
                try:
                    os.remove(temp_path)
                    logger.debug(f"Removed temporary file: {temp_path}")
                except Exception as e:
                    logger.warning(f"Failed to remove temporary file {temp_path}: {str(e)}")
        else:
            # 最適化しない場合は一時ファイルをそのまま使用
            for i, temp_path in enumerate(temp_image_paths):
                # 相対パスを生成
                rel_path = os.path.relpath(temp_path, output_dir)
                image_paths.append(rel_path)
                logger.info(f"Using non-optimized image: {temp_path}")
        
        # PDFファイルを削除
        try:
            os.remove(pdf_path)
            logger.debug(f"Removed temporary PDF file: {pdf_path}")
        except Exception as e:
            logger.warning(f"Failed to remove temporary PDF file: {str(e)}")
        
        logger.info(f"Returning {len(image_paths)} image paths")
        return image_paths, (actual_width, actual_height)
        
    except Exception as e:
        logger.error(f"Exception in convert_to_png: {str(e)}")
        return [], (0, 0)

def convert_coordinates(x: float, y: float, width: float, height: float, 
                      slide_width: float, slide_height: float,
                      image_width: float, image_height: float) -> dict:
    """
    PowerPointの座標系（EMU）をピクセル座標に変換
    
    Args:
        x, y: 元の座標（EMU）
        width, height: 元のサイズ（EMU）
        slide_width: スライドの幅（EMU）
        slide_height: スライドの高さ（EMU）
        image_width: 画像の幅（ピクセル）
        image_height: 画像の高さ（ピクセル）
    
    Returns:
        dict: 変換後の座標とサイズ
    """
    # デバッグ出力
    print(f"Converting coordinates: ({x}, {y}, {width}, {height}) in slide size {slide_width}x{slide_height} to image size {image_width}x{image_height}", file=sys.stderr)
    
    # スライドと画像のアスペクト比を計算
    slide_aspect = slide_width / slide_height
    image_aspect = image_width / image_height
    
    # 基本のスケーリング係数を計算
    scale_x = image_width / slide_width
    scale_y = image_height / slide_height
    
    # アスペクト比の違いを考慮したスケーリング
    if abs(slide_aspect - image_aspect) > 0.01:  # アスペクト比が異なる場合
        print(f"Aspect ratio adjustment needed: slide={slide_aspect}, image={image_aspect}", file=sys.stderr)
        
        # 画像がスライドより横長の場合
        if image_aspect > slide_aspect:
            # 高さに合わせてスケーリング
            new_width = image_height * slide_aspect
            # 中央揃えのためのオフセット計算
            offset_x = (image_width - new_width) / 2
            # スケール係数を調整
            scale_x = new_width / slide_width
            # 座標を変換
            scaled_x = x * scale_x + offset_x
            scaled_y = y * scale_y
        else:
            # 幅に合わせてスケーリング
            new_height = image_width / slide_aspect
            # 中央揃えのためのオフセット計算
            offset_y = (image_height - new_height) / 2
            # スケール係数を調整
            scale_y = new_height / slide_height
            # 座標を変換
            scaled_x = x * scale_x
            scaled_y = y * scale_y + offset_y
    else:
        # アスペクト比が同じ場合は単純にスケーリング
        scaled_x = x * scale_x
        scaled_y = y * scale_y
    
    # サイズをスケーリング
    scaled_width = width * scale_x
    scaled_height = height * scale_y
    
    # デバッグ出力
    print(f"Scaled coordinates: ({scaled_x}, {scaled_y}, {scaled_width}, {scaled_height})", file=sys.stderr)
    
    # 結果を返す
    result = {
        "x": round(scaled_x),
        "y": round(scaled_y),
        "width": round(scaled_width),
        "height": round(scaled_height)
    }
    
    print(f"Final coordinates: {result}", file=sys.stderr)
    return result

def extract_text_from_shape(shape, parent_z_order=0, parent_position=None):
    """シェイプからテキストを抽出し、フォーマット情報も保持する
    
    Args:
        shape: シェイプオブジェクト
        parent_z_order: 親要素のZ順序（グループ内の要素用）
        parent_position: 親要素の位置情報（グループ内の要素用）
        
    Returns:
        dict: テキスト情報を含む辞書、またはNone
    """
    # グループシェイプの場合は再帰的に処理
    if isinstance(shape, GroupShape):
        group_elements = []
        try:
            # グループ自体の位置情報を取得
            group_position = {
                'x': shape.left,
                'y': shape.top,
                'z_order': parent_z_order
            }
            
            # グループ内の各シェイプを処理
            for i, child in enumerate(shape.shapes):
                # 子要素のZ順序は親のZ順序 + インデックスで計算
                child_z_order = parent_z_order + i / 1000.0
                child_element = extract_text_from_shape(child, child_z_order, group_position)
                if child_element:
                    group_elements.append(child_element)
                    
            if group_elements:
                return {
                    "type": "group",
                    "elements": group_elements,
                    "position": group_position
                }
            return None
        except Exception as e:
            logger.warning(f"Error extracting text from group shape: {str(e)}")
            return None
    
    # テーブルの場合は特別な処理
    if hasattr(shape, "table"):
        try:
            table_data = {
                "type": "table",
                "rows": [],
                "position": {
                    'x': shape.left,
                    'y': shape.top,
                    'z_order': parent_z_order
                }
            }
            
            # 各セルのテキストを抽出
            for i, row in enumerate(shape.table.rows):
                row_data = []
                for j, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    if cell_text:
                        # セル内のテキスト書式情報も抽出
                        cell_paragraphs = []
                        for p in cell.text_frame.paragraphs:
                            if not p.text.strip():
                                continue
                                
                            # 各ランのフォント情報を収集
                            runs_info = []
                            for run in p.runs:
                                if not run.text.strip():
                                    continue
                                    
                                font_info = extract_font_info(run)
                                runs_info.append({
                                    "text": run.text.strip(),
                                    "font": font_info,
                                    "index": run._rId
                                })
                                
                            if runs_info:
                                cell_paragraphs.append({
                                    "text": p.text.strip(),
                                    "style": extract_text_style(p),
                                    "runs": runs_info
                                })
                                
                        row_data.append({
                            "text": cell_text,
                            "paragraphs": cell_paragraphs,
                            "position": {
                                "row": i,
                                "column": j
                            }
                        })
                    else:
                        row_data.append(None)  # 空のセル
                        
                table_data["rows"].append(row_data)
                
            return table_data
        except Exception as e:
            logger.warning(f"Error extracting text from table: {str(e)}")
            # テーブル処理に失敗しても通常のテキスト抽出を試みる
    
    # 通常のテキスト抽出処理
    if not hasattr(shape, "text"):
        return None
    
    text = shape.text.strip()
    if not text:
        return None
    
    # シェイプのID、位置情報、サイズを取得
    shape_id = getattr(shape, "shape_id", None)
    left = getattr(shape, "left", 0)
    top = getattr(shape, "top", 0)
    width = getattr(shape, "width", 0)
    height = getattr(shape, "height", 0)
    
    # 親要素の位置がある場合は相対位置を計算
    if parent_position:
        left += parent_position['x']
        top += parent_position['y']
    
    # テキスト情報を抽出
    text_info = {
        "text": text,
        "type": "text",
        "shape_id": shape_id,
        "position": {
            "x": left,
            "y": top,
            "width": width,
            "height": height,
            "z_order": parent_z_order
        }
    }
    
    # 可能であればフォント情報も抽出
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            # テキストフレームの自動調整設定を取得
            if hasattr(shape.text_frame, "auto_size"):
                text_info["auto_size"] = str(shape.text_frame.auto_size)
                
            # 縦書きかどうかを判定
            if hasattr(shape.text_frame, "vertical"):
                text_info["vertical"] = shape.text_frame.vertical
            
            # パラグラフごとの処理
            paragraphs_info = []
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if not p.text.strip():
                    continue
                
                # パラグラフのスタイル情報を取得
                paragraph_style = extract_text_style(p)
                
                # 箇条書きスタイルを取得
                bullet_style = extract_bullet_style(p)
                if bullet_style:
                    paragraph_style.update(bullet_style)
                
                # 各テキストランの処理
                runs_info = []
                for r_idx, run in enumerate(p.runs):
                    if not run.text.strip():
                        continue
                        
                    font_info = extract_font_info(run)
                    runs_info.append({
                        "text": run.text.strip(),
                        "font": font_info,
                        "index": r_idx
                    })
                
                paragraphs_info.append({
                    "text": p.text.strip(),
                    "style": paragraph_style,
                    "runs": runs_info,
                    "index": p_idx
                })
            
            if paragraphs_info:
                text_info["paragraphs"] = paragraphs_info
    except Exception as e:
        # フォント情報の抽出に失敗しても続行
        logger.warning(f"Error extracting font info: {str(e)}")
        
    return text_info

def extract_metadata(presentation: Presentation) -> dict:
    """
    プレゼンテーションのメタデータを抽出する
    
    Args:
        presentation (Presentation): PowerPointプレゼンテーション
    
    Returns:
        dict: メタデータ情報
    """
    core_props = presentation.core_properties
    
    # 安全に属性にアクセスするヘルパー関数
    def safe_get_attr(obj, attr_name, default=None):
        try:
            value = getattr(obj, attr_name)
            if attr_name in ['created', 'modified'] and value:
                return value.isoformat()
            return value or default
        except (AttributeError, TypeError):
            return default
    
    metadata = {
        'title': safe_get_attr(core_props, 'title', ''),
        'author': safe_get_attr(core_props, 'author', ''),
        'created': safe_get_attr(core_props, 'created', ''),
        'modified': safe_get_attr(core_props, 'modified', ''),
        'company': safe_get_attr(core_props, 'company', ''),
        'version': safe_get_attr(core_props, 'version', ''),
        'lastModifiedBy': safe_get_attr(core_props, 'last_modified_by', ''),
        'revision': safe_get_attr(core_props, 'revision', 0),
        'subject': safe_get_attr(core_props, 'subject', ''),
        'keywords': safe_get_attr(core_props, 'keywords', '').split(',') if safe_get_attr(core_props, 'keywords') else [],
        'category': safe_get_attr(core_props, 'category', ''),
        'description': safe_get_attr(core_props, 'description', ''),
        'language': safe_get_attr(core_props, 'language', 'ja-JP'),
        'presentationFormat': 'widescreen' if presentation.slide_width > 9144000 else 'standard',
        'createdApplication': safe_get_attr(core_props, 'application', '')
    }
    
    return metadata

def extract_font_info(run):
    """
    テキストのフォント情報を抽出する
    
    Args:
        run: テキストラン
    
    Returns:
        dict: フォント情報
    """
    font_info = {}
    
    try:
        if run.font:
            font = run.font
            if font.size:
                # EMUからポイントに変換
                font_size = font.size / 12700
                font_info["size"] = font_size
                
            if font.name:
                font_info["name"] = font.name
                
            # フォントファミリー情報を追加
            if hasattr(font, "_Font__family") and font._Font__family:
                font_info["family"] = str(font._Font__family)
                
            # フォントの言語情報を追加
            if hasattr(font, "_Font__charset") and font._Font__charset is not None:
                font_info["charset"] = str(font._Font__charset)
                
            # 文字の計算方向を追加
            if hasattr(font, "_Font__bidi") and font._Font__bidi is not None:
                font_info["bidi"] = font._Font__bidi  # 右から左への書字方向
                
            # スタイル情報を追加
            if font.bold is not None:
                font_info["bold"] = font.bold
                
            if font.italic is not None:
                font_info["italic"] = font.italic
                
            if font.underline is not None:
                font_info["underline"] = font.underline
                
            if hasattr(font, "strike") and font.strike is not None:
                font_info["strike"] = font.strike
                
            # 色情報を追加
            if hasattr(font, "color") and font.color and hasattr(font.color, "rgb"):
                font_info["color"] = safe_get_rgb(font.color)
                
            # ハイライト色情報を追加
            if hasattr(font, "highlight_color") and font.highlight_color:
                font_info["highlight_color"] = str(font.highlight_color)
    except Exception as e:
        logger.warning(f"Error extracting font info: {str(e)}")
        
    return font_info

def extract_text_style(paragraph):
    """
    段落のスタイル情報を抽出する
    
    Args:
        paragraph: 段落オブジェクト
    
    Returns:
        dict: スタイル情報
    """
    style_info = {}
    
    try:
        # 段落の配置
        if paragraph.alignment:
            style_info["alignment"] = str(paragraph.alignment)
            
        # 行間
        if paragraph.line_spacing:
            style_info["line_spacing"] = paragraph.line_spacing
            
        # 段落前の空き
        if paragraph.space_before:
            style_info["space_before"] = paragraph.space_before
            
        # 段落後の空き
        if paragraph.space_after:
            style_info["space_after"] = paragraph.space_after
            
        # インデント情報
        if hasattr(paragraph, "level"):
            style_info["level"] = paragraph.level
            
        # 左インデント
        if hasattr(paragraph, "left_indent") and paragraph.left_indent is not None:
            style_info["left_indent"] = paragraph.left_indent
            
        # 右インデント
        if hasattr(paragraph, "right_indent") and paragraph.right_indent is not None:
            style_info["right_indent"] = paragraph.right_indent
            
        # 最初の行のインデント
        if hasattr(paragraph, "first_line_indent") and paragraph.first_line_indent is not None:
            style_info["first_line_indent"] = paragraph.first_line_indent
            
        # 文字方向（右から左、左から右）
        if hasattr(paragraph, "bidi") and paragraph.bidi is not None:
            style_info["bidi"] = paragraph.bidi
            
        # 縦書きかどうか
        if hasattr(paragraph, "vertical") and paragraph.vertical is not None:
            style_info["vertical"] = paragraph.vertical
            
        # フォントの基本情報を取得（最初のランから）
        if paragraph.runs and paragraph.runs[0].font:
            base_font = paragraph.runs[0].font
            if base_font.name:
                style_info["font_name"] = base_font.name
                
            if base_font.size:
                style_info["font_size"] = base_font.size / 12700  # EMUからポイントに変換
    except Exception as e:
        logger.warning(f"Error extracting text style: {str(e)}")
        
    return style_info

def extract_bullet_style(paragraph):
    """
    箇条書きのスタイル情報を抽出する
    
    Args:
        paragraph: 段落オブジェクト
    
    Returns:
        dict: 箇条書きスタイル情報
    """
    bullet_info = {}
    
    try:
        # 箇条書きの有無を確認
        has_bullet = False
        if hasattr(paragraph, "bullet") and paragraph.bullet:
            has_bullet = True
        elif hasattr(paragraph, "level") and paragraph.level is not None and paragraph.level > 0:
            # レベルが設定されている場合も箇条書きの可能性がある
            has_bullet = True
            
        if not has_bullet:
            return None
            
        bullet_info["has_bullet"] = True
        
        # 箇条書きのレベル
        if hasattr(paragraph, "level"):
            bullet_info["level"] = paragraph.level
        
        # 箇条書きの詳細情報を取得
        if hasattr(paragraph, "bullet") and paragraph.bullet:
            bullet = paragraph.bullet
            
            # 箇条書きの種類
            bullet_info["type"] = 'number' if bullet.numbered else 'bullet'
            
            # 箇条書きの書式
            if hasattr(bullet, "format"):
                bullet_info["format"] = bullet.format
                
            # 開始番号
            if hasattr(bullet, "start_num"):
                bullet_info["start_at"] = bullet.start_num
            
            # 箇条書きの文字
            if hasattr(bullet, "char") and bullet.char:
                # Unicode正規化を適用して特殊文字を適切に処理
                import unicodedata
                bullet_info["char"] = unicodedata.normalize('NFKC', bullet.char)
            
            # 箇条書きのフォント
            if hasattr(bullet, "font") and bullet.font:
                bullet_info["font"] = bullet.font
            
            # 箇条書きの色
            if hasattr(bullet, "color") and bullet.color:
                bullet_info["color"] = safe_get_rgb(bullet.color)
            
            # 箇条書きのサイズ
            if hasattr(bullet, "size") and bullet.size:
                bullet_info["size"] = bullet.size / 12700  # EMUからポイントに変換
    except Exception as e:
        logger.warning(f"Error extracting bullet style: {str(e)}")
        return None
        
    return bullet_info

def extract_shape_info(shape) -> dict:
    """
    シェイプの情報を抽出する
    
    Args:
        shape: シェイプオブジェクト
    
    Returns:
        dict: シェイプ情報
    """
    shape_type = shape.shape_type
    
    base_info = {
        'type': str(shape_type),
        'x': shape.left,
        'y': shape.top,
        'width': shape.width,
        'height': shape.height,
        'rotation': shape.rotation,
        'fillColor': extract_fill_color(shape),
        'strokeColor': extract_line_color(shape),
        'strokeWidth': shape.line.width if shape.line else 1,
        'opacity': 1  # python-pptxの最新バージョンではtransparency属性が存在しないため、固定値を使用
    }
    
    # 円の場合は半径を追加
    if shape_type == 'OVAL':
        base_info['radius'] = min(shape.width, shape.height) / 2
        
    # 多角形の場合は頂点情報を追加
    if hasattr(shape, 'points'):
        base_info['points'] = [{'x': point[0], 'y': point[1]} for point in shape.points]
        
    return base_info

def safe_get_rgb(color) -> str:
    """
    色情報からRGB値を安全に取得する
    
    Args:
        color: 色オブジェクト
        
    Returns:
        str: 16進数のRGB値
    """
    if not color:
        return '#000000'
        
    try:
        if hasattr(color, 'rgb') and color.rgb:
            return f'#{color.rgb:06x}'
        elif hasattr(color, 'theme_color') and color.theme_color:
            # テーマカラーの場合
            return '#000000'  # デフォルト値
    except (AttributeError, TypeError):
        pass
        
    return '#000000'

def extract_fill_color(shape) -> str:
    """
    シェイプの塗りつぶし色を抽出する
    
    Args:
        shape: シェイプオブジェクト
    
    Returns:
        str: 色情報（16進数）
    """
    if not hasattr(shape, 'fill') or not shape.fill or shape.fill.type == 'NONE':
        return 'transparent'
        
    if shape.fill.type == 'SOLID':
        return safe_get_rgb(shape.fill.fore_color)
        
    return 'transparent'

def extract_line_color(shape) -> str:
    """
    シェイプの線の色を抽出する
    
    Args:
        shape: シェイプオブジェクト
    
    Returns:
        str: 色情報（16進数）
    """
    try:
        if not hasattr(shape, 'line') or not shape.line:
            return 'transparent'
            
        if not hasattr(shape.line, 'fill') or not shape.line.fill or shape.line.fill.type == 'NONE':
            return 'transparent'
            
        if hasattr(shape.line, 'color') and shape.line.color:
            return safe_get_rgb(shape.line.color)
    except (AttributeError, TypeError):
        pass
        
    return 'transparent'

def extract_background(slide) -> dict:
    """
    スライドの背景情報を抽出する
    
    Args:
        slide: スライドオブジェクト
    
    Returns:
        dict: 背景情報
    """
    background = slide.background
    
    if not background:
        return {
            'color': '#FFFFFF',
            'image': None,
            'pattern': None,
            'gradient': None,
            'transparency': 0
        }
        
    fill = background.fill
    
    # python-pptxの最新バージョンではtransparency属性が存在しないため、固定値を使用
    transparency = 0
    
    result = {
        'color': '#FFFFFF',
        'image': None,
        'pattern': None,
        'gradient': None,
        'transparency': transparency
    }
    
    if fill:
        try:
            if fill.type == 'SOLID':
                result['color'] = safe_get_rgb(fill.fore_color) if hasattr(fill, 'fore_color') else '#FFFFFF'
            elif fill.type == 'PICTURE' and hasattr(fill, 'image'):
                try:
                    result['image'] = {
                        'rId': fill.image.rId if hasattr(fill.image, 'rId') else '',
                        'filename': fill.image.filename if hasattr(fill.image, 'filename') else ''
                    }
                except (AttributeError, TypeError):
                    result['image'] = {'rId': '', 'filename': ''}
            elif fill.type == 'PATTERN' and hasattr(fill, 'pattern'):
                result['pattern'] = {
                    'type': str(fill.pattern),
                    'foreColor': safe_get_rgb(fill.fore_color) if hasattr(fill, 'fore_color') else '#000000',
                    'backColor': safe_get_rgb(fill.back_color) if hasattr(fill, 'back_color') else '#FFFFFF'
                }
        except (AttributeError, TypeError) as e:
            print(f"Error processing fill: {str(e)}", file=sys.stderr)
            
        # グラデーションの処理
        try:
            if fill.type == 'GRADIENT' and hasattr(fill, 'gradient_stops'):
                stops = []
                for stop in fill.gradient_stops:
                    try:
                        stops.append({
                            'position': stop.position if hasattr(stop, 'position') else 0,
                            'color': safe_get_rgb(stop.color) if hasattr(stop, 'color') else '#000000'
                        })
                    except (AttributeError, TypeError) as e:
                        print(f"Error processing gradient stop: {str(e)}", file=sys.stderr)
                
                result['gradient'] = {
                    'type': 'linear' if fill.gradient_stops else 'radial',
                    'stops': stops,
                    'angle': fill.gradient_angle if hasattr(fill, 'gradient_angle') else 0
                }
        except (AttributeError, TypeError) as e:
            print(f"Error processing gradient: {str(e)}", file=sys.stderr)
            
    return result

def parse_pptx(file_path: str, output_dir: str, optimize_images: bool = True, 
             image_format: str = 'WEBP', image_quality: int = 85,
             max_width: int = 1920, max_height: int = 1080,
             sort_elements: bool = True, extract_tables: bool = True,
             extract_groups: bool = True, extract_smartart: bool = True,
             improve_text_order: bool = True) -> Dict[str, Any]:
    """
    PPTXファイルを解析し、スライド情報を抽出する
    
    Args:
        file_path (str): PPTXファイルのパス
        output_dir (str): 出力ディレクトリ
        optimize_images (bool): 画像最適化を行うかどうか
        image_format (str): 出力画像形式 ('PNG', 'WEBP', 'JPEG')
        image_quality (int): 画質 (0-100)
        max_width (int): 最大幅
        max_height (int): 最大高さ
        sort_elements (bool): 要素をZ順序とレイアウト位置でソートするか
        extract_tables (bool): テーブル内のテキストを抽出するか
        extract_groups (bool): グループ内のテキストを抽出するか
        extract_smartart (bool): SmartArt内のテキストを抽出するか
        improve_text_order (bool): テキスト順序を改善するか
        
    Returns:
        Dict[str, Any]: 解析結果
    """
    try:
        logger.info(f"Starting to parse PPTX file: {file_path}")
        # プレゼンテーションを開く
        presentation = Presentation(file_path)
        
        # スライド画像を生成（最適化オプション付き）
        image_paths, image_size = convert_to_png(
            file_path, 
            output_dir,
            optimize=optimize_images,
            format=image_format,
            quality=image_quality,
            max_width=max_width,
            max_height=max_height
        )
        
        if not image_paths:
            logger.error("Failed to convert slides to images")
            return {'error': 'Failed to convert slides to images'}
        
        # 結果を格納する辞書
        result = {
            'slides': [],
            'metadata': {
                'image_format': image_format if optimize_images else 'PNG',
                'image_quality': image_quality,
                'image_dimensions': {
                    'width': image_size[0],
                    'height': image_size[1]
                },
                'optimized': optimize_images,
                'parser_version': '2.0',  # パーサーバージョンを追加
                'extraction_options': {
                    'tables': extract_tables,
                    'groups': extract_groups,
                    'smartart': extract_smartart,
                    'improved_text_order': improve_text_order
                }
            }
        }
        
        # メタデータを抽出
        try:
            result['metadata'].update(extract_metadata(presentation))
        except Exception as e:
            logger.warning(f"Error extracting metadata: {str(e)}")
        
        # 各スライドを処理
        total_slides = len(presentation.slides)
        logger.info(f"Processing {total_slides} slides")
        
        for i, (slide, image_path) in enumerate(zip(presentation.slides, image_paths)):
            logger.info(f"Processing slide {i+1}/{total_slides}")
            
            slide_data = {
                'index': i,
                'image_path': image_path,
                'elements': [],
                'text_elements': [],  # テキスト要素のみを格納する配列を追加
                'background': extract_background(slide),
                'size': {
                    'width': presentation.slide_width,
                    'height': presentation.slide_height
                }
            }
            
            # スライド内の全要素を処理
            all_elements = []
            for shape_idx, shape in enumerate(slide.shapes):
                # Z順序を設定（スライド内の順序を使用）
                z_order = shape_idx
                
                # テキスト要素の抽出
                if hasattr(shape, "text") or isinstance(shape, GroupShape) or hasattr(shape, "table"):
                    text_data = extract_text_from_shape(shape, z_order)
                    if text_data:
                        all_elements.append(text_data)
                        # テキスト要素のみの配列にも追加
                        if text_data.get("type") == "text" or text_data.get("type") == "table" or text_data.get("type") == "group":
                            slide_data['text_elements'].append(text_data)
                
                # 非テキスト要素の情報も抽出
                shape_info = extract_shape_info(shape)
                if shape_info:
                    all_elements.append(shape_info)
            
            # 要素をZ順序とレイアウト位置でソート
            if sort_elements and all_elements:
                # まずZ順序でソート
                all_elements.sort(key=lambda x: x.get("position", {}).get("z_order", 0) if x.get("position") else 999999)
                
                # 読み順を改善する場合は、さらにレイアウト位置も考慮
                if improve_text_order:
                    # テキスト要素のみを抽出してレイアウト位置でソート
                    text_elements = [elem for elem in all_elements if elem.get("type") in ["text", "table", "group"]]
                    
                    # 縦書きテキストと横書きテキストを区別
                    vertical_texts = []
                    horizontal_texts = []
                    other_elements = []
                    
                    for elem in text_elements:
                        if elem.get("type") == "text" and elem.get("vertical"):
                            vertical_texts.append(elem)
                        elif elem.get("type") == "text":
                            horizontal_texts.append(elem)
                        else:
                            other_elements.append(elem)
                    
                    # 横書きテキストは上から下、左から右の順でソート
                    horizontal_texts.sort(key=lambda x: (x.get("position", {}).get("y", 0), x.get("position", {}).get("x", 0)))
                    
                    # 縦書きテキストは右から左、上から下の順でソート
                    vertical_texts.sort(key=lambda x: (-x.get("position", {}).get("x", 0), x.get("position", {}).get("y", 0)))
                    
                    # ソートされたテキスト要素を結合
                    sorted_text_elements = horizontal_texts + vertical_texts + other_elements
                    
                    # テキスト要素のみの配列を更新
                    slide_data['text_elements'] = sorted_text_elements
                    
                    # 全要素の配列でテキスト要素を更新
                    text_element_ids = {id(elem) for elem in sorted_text_elements}
                    non_text_elements = [elem for elem in all_elements if id(elem) not in text_element_ids]
                    all_elements = sorted_text_elements + non_text_elements
            
            # 処理済みの要素をスライドデータに追加
            slide_data['elements'] = all_elements
            
            # 特殊文字や多言語テキストの処理を強化
            for elem in slide_data['text_elements']:
                if elem.get("type") == "text" and "text" in elem:
                    # Unicode正規化を適用して特殊文字を適切に処理
                    import unicodedata
                    normalized_text = unicodedata.normalize('NFKC', elem["text"])
                    elem["text"] = normalized_text
                    
                    # パラグラフ内のテキストも正規化
                    if "paragraphs" in elem:
                        for para in elem["paragraphs"]:
                            if "text" in para:
                                para["text"] = unicodedata.normalize('NFKC', para["text"])
                            
                            # ラン内のテキストも正規化
                            if "runs" in para:
                                for run in para["runs"]:
                                    if "text" in run:
                                        run["text"] = unicodedata.normalize('NFKC', run["text"])
            
            # スライドデータを結果に追加
            result['slides'].append(slide_data)
            
        logger.info(f"Successfully processed all {total_slides} slides")
        return result
        
    except Exception as e:
        logger.error(f"Error parsing PPTX: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        return {'error': str(e), 'traceback': traceback.format_exc()}

def update_pptx_with_translations(
    original_file: str,
    output_file: str,
    translations: Dict[int, List[Dict[str, str]]]
) -> bool:
    """
    翻訳されたテキストでPPTXファイルを更新する
    
    Args:
        original_file (str): 元のPPTXファイルパス
        output_file (str): 出力するPPTXファイルパス
        translations (Dict[int, List[Dict[str, str]]]): 
            スライドインデックスごとの翻訳テキスト情報
    
    Returns:
        bool: 成功した場合True
    """
    try:
        prs = Presentation(original_file)
        
        for slide_index, slide in enumerate(prs.slides):
            if slide_index not in translations:
                continue
                
            slide_translations = translations[slide_index]
            shape_map = {shape.shape_id: shape for shape in slide.shapes 
                        if hasattr(shape, "text")}
            
            for trans in slide_translations:
                shape_id = trans.get("shape_id")
                if shape_id in shape_map:
                    shape = shape_map[shape_id]
                    shape.text = trans["translated_text"]
        
        prs.save(output_file)
        return True
        
    except Exception:
        return False

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='PPTX Parser')
    parser.add_argument('file_path', help='Path to PPTX file')
    parser.add_argument('output_dir', help='Output directory')
    
    # 画像関連のオプション
    image_group = parser.add_argument_group('Image options')
    image_group.add_argument('--optimize', action='store_true', default=True, help='Optimize images')
    image_group.add_argument('--no-optimize', action='store_false', dest='optimize', help='Disable image optimization')
    image_group.add_argument('--format', choices=['PNG', 'WEBP', 'JPEG'], default='WEBP', help='Image format')
    image_group.add_argument('--quality', type=int, default=85, help='Image quality (0-100)')
    image_group.add_argument('--max-width', type=int, default=1920, help='Maximum image width')
    image_group.add_argument('--max-height', type=int, default=1080, help='Maximum image height')
    
    # テキスト抽出関連のオプション
    text_group = parser.add_argument_group('Text extraction options')
    text_group.add_argument('--sort-elements', action='store_true', default=True, 
                          help='Sort elements by Z-order and layout position')
    text_group.add_argument('--no-sort-elements', action='store_false', dest='sort_elements', 
                          help='Do not sort elements')
    text_group.add_argument('--extract-tables', action='store_true', default=True, 
                          help='Extract text from tables')
    text_group.add_argument('--no-extract-tables', action='store_false', dest='extract_tables', 
                          help='Do not extract text from tables')
    text_group.add_argument('--extract-groups', action='store_true', default=True, 
                          help='Extract text from grouped objects')
    text_group.add_argument('--no-extract-groups', action='store_false', dest='extract_groups', 
                          help='Do not extract text from grouped objects')
    text_group.add_argument('--extract-smartart', action='store_true', default=True, 
                          help='Extract text from SmartArt')
    text_group.add_argument('--no-extract-smartart', action='store_false', dest='extract_smartart', 
                          help='Do not extract text from SmartArt')
    text_group.add_argument('--improve-text-order', action='store_true', default=True, 
                          help='Improve text reading order')
    text_group.add_argument('--no-improve-text-order', action='store_false', dest='improve_text_order', 
                          help='Do not improve text reading order')
    
    # ログ関連のオプション
    log_group = parser.add_argument_group('Logging options')
    log_group.add_argument('--verbose', '-v', action='store_true', help='Enable verbose logging')
    log_group.add_argument('--debug', action='store_true', help='Enable debug logging')
    
    args = parser.parse_args()
    
    # ログレベルの設定
    if args.debug:
        logger.setLevel(logging.DEBUG)
    elif args.verbose:
        logger.setLevel(logging.INFO)
    else:
        logger.setLevel(logging.WARNING)
    
    # 入力ファイルが存在するか確認
    if not os.path.exists(args.file_path):
        logger.error(f"Input file '{args.file_path}' does not exist")
        print(f"Error: Input file '{args.file_path}' does not exist", file=sys.stderr)
        sys.exit(1)
        
    # 出力ディレクトリが存在しない場合は作成
    os.makedirs(args.output_dir, exist_ok=True)
    
    logger.info(f"Starting to parse PPTX file: {args.file_path}")
    logger.info(f"Output directory: {args.output_dir}")
    logger.info(f"Image options: format={args.format}, quality={args.quality}, optimize={args.optimize}")
    logger.info(f"Text extraction options: sort_elements={args.sort_elements}, extract_tables={args.extract_tables}, " +
               f"extract_groups={args.extract_groups}, extract_smartart={args.extract_smartart}, " +
               f"improve_text_order={args.improve_text_order}")
    
    # PPTXファイルを解析
    result = parse_pptx(
        args.file_path, 
        args.output_dir,
        optimize_images=args.optimize,
        image_format=args.format,
        image_quality=args.quality,
        max_width=args.max_width,
        max_height=args.max_height,
        sort_elements=args.sort_elements,
        extract_tables=args.extract_tables,
        extract_groups=args.extract_groups,
        extract_smartart=args.extract_smartart,
        improve_text_order=args.improve_text_order
    )
    
    # 結果を出力
    print(json.dumps(result, ensure_ascii=False, indent=2))
    
    # 処理結果のサマリをログに出力
    total_slides = len(result.get('slides', []))
    total_text_elements = sum(len(slide.get('text_elements', [])) for slide in result.get('slides', []))
    logger.info(f"Successfully processed {total_slides} slides with {total_text_elements} text elements")
    
    # エラーがあればログに出力
    if 'error' in result:
        logger.error(f"Error in parsing: {result['error']}")
        sys.exit(1)
