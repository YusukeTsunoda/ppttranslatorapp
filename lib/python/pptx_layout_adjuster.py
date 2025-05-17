#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPTXレイアウト調整モジュール
翻訳後のPPTXファイルのレイアウトを調整するための機能を提供します。
"""

import os
import sys
import json
import logging
import tempfile
import shutil
from typing import Dict, List, Any, Tuple, Optional, Union
import math
import re
from datetime import datetime

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table, _Cell
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt, Inches, Cm

# 自作モジュールをインポート
from layout_utils import (
    TextBoxInfo, AdjustedTextBoxInfo, adjust_layout,
    get_expansion_ratio, LANGUAGE_OFFSETS
)

# ロギング設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('pptx_layout_adjuster')

class PPTXLayoutAdjuster:
    """PPTXファイルのレイアウトを調整するクラス"""
    
    def __init__(self, options: Dict[str, Any] = None):
        """
        コンストラクタ
        
        Args:
            options: レイアウト調整オプション
        """
        # デフォルトオプション
        self.default_options = {
            'enable_auto_resize': True,
            'resize_strategy': 'both',
            'max_width_expansion': 1.5,
            'max_height_expansion': 1.5,
            'apply_font_mapping': True,
            'apply_language_offset': True,
            'detect_collisions': True,
            'safety_margin': 0.1,
            'preserve_original_file': True,
            'generate_preview_images': False,
            'apply_font_embedding': False,
            'generate_report': True
        }
        
        # オプションをマージ
        self.options = self.default_options.copy()
        if options:
            self.options.update(options)
        
        # 結果を格納する変数
        self.result = {
            'success': False,
            'adjusted_text_boxes': [],
            'warnings': [],
            'errors': [],
            'overflow_count': 0,
            'collision_count': 0
        }
    
    def adjust_pptx_layout(
        self,
        input_file_path: str,
        output_file_path: str,
        translations: List[Dict[str, Any]],
        source_language: str,
        target_language: str
    ) -> Dict[str, Any]:
        """
        PPTXファイルのレイアウトを調整する
        
        Args:
            input_file_path: 入力ファイルパス
            output_file_path: 出力ファイルパス
            translations: 翻訳情報のリスト
            source_language: 翻訳元言語コード
            target_language: 翻訳先言語コード
            
        Returns:
            レイアウト調整結果
        """
        # 結果オブジェクトを初期化
        self.result = {
            'success': False,
            'adjusted_text_boxes': [],
            'warnings': [],
            'errors': [],
            'overflow_count': 0,
            'collision_count': 0
        }
        
        try:
            # 入力ファイルの存在確認
            if not os.path.exists(input_file_path):
                self.result['errors'].append(f"入力ファイルが見つかりません: {input_file_path}")
                return self.result
            
            # 出力ディレクトリの作成
            output_dir = os.path.dirname(output_file_path)
            if not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # 元のファイルを保存
            if self.options['preserve_original_file']:
                original_file_path = os.path.join(
                    output_dir,
                    f"{os.path.splitext(os.path.basename(output_file_path))[0]}_original.pptx"
                )
                shutil.copy2(input_file_path, original_file_path)
                logger.info(f"元のファイルを保存しました: {original_file_path}")
            
            # 翻訳情報をシェイプIDでインデックス化
            translation_map = {}
            for translation in translations:
                slide_index = translation.get('slideIndex', 0)
                shape_id = translation.get('shapeId', '')
                
                if not shape_id:
                    continue
                
                key = f"{slide_index}_{shape_id}"
                translation_map[key] = translation
            
            # PPTXファイルを読み込む
            logger.info(f"PPTXファイルを読み込んでいます: {input_file_path}")
            prs = Presentation(input_file_path)
            
            # 各スライドを処理
            for slide_index, slide in enumerate(prs.slides):
                logger.info(f"スライド {slide_index + 1} を処理しています...")
                
                # スライド内のテキストボックス情報を収集
                text_boxes = []
                shape_map = {}
                
                # シェイプを処理
                for shape_index, shape in enumerate(slide.shapes):
                    shape_id = str(shape_index)
                    
                    # テキストボックス情報を抽出
                    text_box_info = self._extract_text_box_info(
                        shape, shape_id, slide_index
                    )
                    
                    if text_box_info:
                        # 翻訳情報があれば適用
                        key = f"{slide_index}_{shape_id}"
                        if key in translation_map:
                            translation = translation_map[key]
                            text_box_info.text = translation.get('translatedText', text_box_info.text)
                        
                        text_boxes.append(text_box_info)
                        shape_map[shape_id] = shape
                
                # レイアウトを調整
                if text_boxes:
                    adjusted_text_boxes = adjust_layout(
                        text_boxes,
                        source_language,
                        target_language,
                        self.options
                    )
                    
                    # 調整結果をPPTXに適用
                    for adjusted in adjusted_text_boxes:
                        shape_id = adjusted.id
                        if shape_id in shape_map:
                            shape = shape_map[shape_id]
                            self._apply_adjustment_to_shape(shape, adjusted)
                            
                            # 結果を記録
                            self.result['adjusted_text_boxes'].append(adjusted.to_dict())
                            
                            # オーバーフローと衝突をカウント
                            if adjusted.overflow_detected:
                                self.result['overflow_count'] += 1
                            
                            if adjusted.collision_detected:
                                self.result['collision_count'] += 1
            
            # PPTXファイルを保存
            logger.info(f"PPTXファイルを保存しています: {output_file_path}")
            prs.save(output_file_path)
            
            # 警告メッセージを生成
            if self.result['overflow_count'] > 0:
                self.result['warnings'].append(
                    f"{self.result['overflow_count']}件のテキストボックスでオーバーフローを検出しました"
                )
            
            if self.result['collision_count'] > 0:
                self.result['warnings'].append(
                    f"{self.result['collision_count']}件のテキストボックスで衝突を検出しました"
                )
            
            # レポートを生成
            if self.options['generate_report']:
                report_path = os.path.join(
                    output_dir,
                    f"{os.path.splitext(os.path.basename(output_file_path))[0]}_layout_report.json"
                )
                
                with open(report_path, 'w', encoding='utf-8') as f:
                    json.dump({
                        'input_file': input_file_path,
                        'output_file': output_file_path,
                        'source_language': source_language,
                        'target_language': target_language,
                        'options': self.options,
                        'adjustment_result': {
                            'total_text_boxes': len(self.result['adjusted_text_boxes']),
                            'overflow_count': self.result['overflow_count'],
                            'collision_count': self.result['collision_count'],
                            'warnings': self.result['warnings'],
                            'errors': self.result['errors']
                        },
                        'timestamp': datetime.now().isoformat()
                    }, f, indent=2, ensure_ascii=False)
                
                self.result['report_path'] = report_path
                logger.info(f"レイアウト調整レポートを生成しました: {report_path}")
            
            # 成功フラグを設定
            self.result['success'] = len(self.result['errors']) == 0
            
            return self.result
            
        except Exception as e:
            error_message = str(e)
            self.result['errors'].append(f"PPTXレイアウト調整中にエラーが発生しました: {error_message}")
            logger.error(f"PPTXレイアウト調整中にエラーが発生しました: {error_message}", exc_info=True)
            
            return self.result
    
    def _extract_text_box_info(
        self,
        shape: Any,
        shape_id: str,
        slide_index: int,
        parent_z_order: int = 0
    ) -> Optional[TextBoxInfo]:
        """
        シェイプからテキストボックス情報を抽出する
        
        Args:
            shape: シェイプオブジェクト
            shape_id: シェイプID
            slide_index: スライドインデックス
            parent_z_order: 親要素のZ順序
            
        Returns:
            テキストボックス情報（テキストがない場合はNone）
        """
        try:
            # グループシェイプの場合は処理しない（個別のシェイプとして処理される）
            if isinstance(shape, GroupShape):
                return None
            
            # テーブルの場合は処理しない（セルごとに処理される）
            if hasattr(shape, 'has_table') and shape.has_table:
                return None
            
            # テキストがない場合は処理しない
            if not hasattr(shape, 'text') or not shape.text.strip():
                return None
            
            # 位置情報を取得
            x = shape.left
            y = shape.top
            width = shape.width
            height = shape.height
            
            # Z順序を取得（親要素のZ順序を考慮）
            z_order = parent_z_order
            
            # フォント情報を取得
            font_name = None
            font_size = None
            font_color = None
            alignment = None
            vertical = False
            rtl = False
            
            # テキストフレームがある場合
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                
                # 縦書きかどうか
                vertical = hasattr(text_frame, 'vertical') and text_frame.vertical
                
                # 段落があれば最初の段落からフォント情報を取得
                if hasattr(text_frame, 'paragraphs') and text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    
                    # 配置を取得
                    if hasattr(paragraph, 'alignment'):
                        alignment_map = {
                            PP_ALIGN.LEFT: 'left',
                            PP_ALIGN.CENTER: 'center',
                            PP_ALIGN.RIGHT: 'right',
                            PP_ALIGN.JUSTIFY: 'justify'
                        }
                        alignment = alignment_map.get(paragraph.alignment, 'left')
                    
                    # RTL（右から左）かどうか
                    rtl = hasattr(paragraph, 'rtl') and paragraph.rtl
                    
                    # ランがあればフォント情報を取得
                    if hasattr(paragraph, 'runs') and paragraph.runs:
                        run = paragraph.runs[0]
                        
                        if hasattr(run, 'font'):
                            font = run.font
                            
                            # フォント名
                            if hasattr(font, 'name') and font.name:
                                font_name = font.name
                            
                            # フォントサイズ
                            if hasattr(font, 'size') and font.size:
                                font_size = font.size.pt
                            
                            # フォントカラー
                            if hasattr(font, 'color') and hasattr(font.color, 'rgb') and font.color.rgb:
                                rgb = font.color.rgb
                                font_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
            
            # テキストボックス情報を作成
            return TextBoxInfo(
                id=shape_id,
                text=shape.text,
                x=x,
                y=y,
                width=width,
                height=height,
                font_name=font_name,
                font_size=font_size,
                font_color=font_color,
                alignment=alignment,
                vertical=vertical,
                rtl=rtl,
                z_order=z_order
            )
            
        except Exception as e:
            logger.warning(f"テキストボックス情報の抽出中にエラーが発生しました: {e}", exc_info=True)
            return None
    
    def _apply_adjustment_to_shape(self, shape: Any, adjusted: AdjustedTextBoxInfo) -> None:
        """
        調整結果をシェイプに適用する
        
        Args:
            shape: シェイプオブジェクト
            adjusted: 調整後のテキストボックス情報
        """
        try:
            # テキストを設定
            if hasattr(shape, 'text'):
                shape.text = adjusted.text
            
            # 位置とサイズを設定
            shape.left = adjusted.x
            shape.top = adjusted.y
            shape.width = adjusted.width
            shape.height = adjusted.height
            
            # テキストフレームがある場合
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                
                # 自動調整を無効化（サイズを固定）
                if hasattr(text_frame, 'word_wrap'):
                    text_frame.word_wrap = True
                
                if hasattr(text_frame, 'auto_size'):
                    text_frame.auto_size = None
                
                # 段落があれば設定を適用
                if hasattr(text_frame, 'paragraphs') and text_frame.paragraphs:
                    for paragraph in text_frame.paragraphs:
                        # 配置を設定
                        if adjusted.alignment and hasattr(paragraph, 'alignment'):
                            alignment_map = {
                                'left': PP_ALIGN.LEFT,
                                'center': PP_ALIGN.CENTER,
                                'right': PP_ALIGN.RIGHT,
                                'justify': PP_ALIGN.JUSTIFY
                            }
                            if adjusted.alignment in alignment_map:
                                paragraph.alignment = alignment_map[adjusted.alignment]
                        
                        # RTL（右から左）を設定
                        if hasattr(paragraph, 'rtl'):
                            paragraph.rtl = adjusted.rtl
                        
                        # ランがあればフォント情報を設定
                        if hasattr(paragraph, 'runs') and paragraph.runs:
                            for run in paragraph.runs:
                                if hasattr(run, 'font'):
                                    font = run.font
                                    
                                    # フォント名
                                    if adjusted.font_name and hasattr(font, 'name'):
                                        font.name = adjusted.font_name
                                    
                                    # フォントサイズ
                                    if adjusted.font_size and hasattr(font, 'size'):
                                        font.size = Pt(adjusted.font_size)
                                    
                                    # フォントカラー
                                    if adjusted.font_color and hasattr(font, 'color') and hasattr(font.color, 'rgb'):
                                        # 16進数の色コードをRGB値に変換
                                        if adjusted.font_color.startswith('#') and len(adjusted.font_color) == 7:
                                            r = int(adjusted.font_color[1:3], 16)
                                            g = int(adjusted.font_color[3:5], 16)
                                            b = int(adjusted.font_color[5:7], 16)
                                            font.color.rgb = (r, g, b)
            
        except Exception as e:
            logger.warning(f"シェイプへの調整適用中にエラーが発生しました: {e}", exc_info=True)

def main():
    """メイン関数"""
    import argparse
    
    # コマンドライン引数の解析
    parser = argparse.ArgumentParser(description='PPTXファイルのレイアウトを調整します')
    parser.add_argument('input_file', help='入力PPTXファイルのパス')
    parser.add_argument('output_file', help='出力PPTXファイルのパス')
    parser.add_argument('translation_file', help='翻訳情報JSONファイルのパス')
    parser.add_argument('--source', default='ja', help='翻訳元言語コード（デフォルト: ja）')
    parser.add_argument('--target', default='en', help='翻訳先言語コード（デフォルト: en）')
    parser.add_argument('--config', help='設定JSONファイルのパス')
    
    args = parser.parse_args()
    
    # 設定ファイルの読み込み
    options = {}
    if args.config and os.path.exists(args.config):
        with open(args.config, 'r', encoding='utf-8') as f:
            options = json.load(f)
    
    # 翻訳情報の読み込み
    translations = []
    if os.path.exists(args.translation_file):
        with open(args.translation_file, 'r', encoding='utf-8') as f:
            translations = json.load(f)
    
    # PPTXレイアウト調整
    adjuster = PPTXLayoutAdjuster(options)
    result = adjuster.adjust_pptx_layout(
        args.input_file,
        args.output_file,
        translations,
        args.source,
        args.target
    )
    
    # 結果を出力
    print(json.dumps(result, indent=2, ensure_ascii=False))
    
    # 成功/失敗を表示
    if result['success']:
        print(f"レイアウト調整が完了しました: {args.output_file}")
        if 'report_path' in result:
            print(f"レポート: {result['report_path']}")
    else:
        print("レイアウト調整中にエラーが発生しました:")
        for error in result['errors']:
            print(f"- {error}")

if __name__ == "__main__":
    main()
