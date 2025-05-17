#!/usr/bin/env python3
import sys
import json
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from typing import Dict, List, Any, Optional, Tuple
import math
import logging

# ロギングの設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('pptx_generator')

# EMUからピクセル、ピクセルからEMUへの変換関数
def emu_to_pixels(emu: int) -> float:
    """EMU単位をピクセルに変換"""
    return emu / 9525  # 1 pixel = 9525 EMU

def pixels_to_emu(pixels: float) -> int:
    """ピクセルをEMU単位に変換"""
    return int(pixels * 9525)  # 1 pixel = 9525 EMU

def apply_text_style(run, style: Dict):
    """テキストスタイルを適用"""
    try:
        if "fontSize" in style:
            run.font.size = Pt(style["fontSize"])
        if "fontName" in style:
            run.font.name = style["fontName"]
        elif "fontFamily" in style:
            run.font.name = style["fontFamily"]
        if "bold" in style:
            run.font.bold = style["bold"]
        if "italic" in style:
            run.font.italic = style["italic"]
        if "color" in style and isinstance(style["color"], dict):
            color = style["color"]
            if all(k in color for k in ['r', 'g', 'b']):
                run.font.color.rgb = RGBColor(color['r'], color['g'], color['b'])
        if "underline" in style:
            run.font.underline = style["underline"]
    except Exception as e:
        logger.warning(f"スタイル適用エラー: {e}")

def get_alignment(alignment_str: str) -> PP_ALIGN:
    """テキスト配置を文字列からPP_ALIGN値に変換"""
    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
        "distribute": PP_ALIGN.DISTRIBUTE
    }
    return alignment_map.get(alignment_str.lower(), PP_ALIGN.LEFT)

def estimate_text_height(text: str, font_size: float, width: float, lang: str = "en") -> float:
    """テキストの高さをより正確に推定"""
    # 言語ごとの文字幅係数
    char_width_factors = {
        "en": 0.6,  # 英語（基準）
        "ja": 1.0,  # 日本語
        "zh": 1.0,  # 中国語
        "ko": 0.9,  # 韓国語
        "ru": 0.7,  # ロシア語
        "ar": 0.8,  # アラビア語
        "de": 0.65, # ドイツ語
        "fr": 0.65, # フランス語
        "es": 0.65, # スペイン語
    }
    
    # 言語ごとの行高係数
    line_height_factors = {
        "en": 1.2,  # 英語（基準）
        "ja": 1.5,  # 日本語
        "zh": 1.5,  # 中国語
        "ko": 1.4,  # 韓国語
        "ru": 1.3,  # ロシア語
        "ar": 1.4,  # アラビア語
        "de": 1.2,  # ドイツ語
        "fr": 1.2,  # フランス語
        "es": 1.2,  # スペイン語
    }
    
    # 言語に応じた係数を取得
    char_width_factor = char_width_factors.get(lang, 0.6)
    line_height_factor = line_height_factors.get(lang, 1.2)
    
    # 改行文字で分割して行数を計算
    explicit_lines = text.split('\n')
    total_lines = 0
    
    for line in explicit_lines:
        if not line:  # 空行の場合
            total_lines += 1
            continue
            
        # 各行の文字数から必要な行数を計算
        avg_char_width = font_size * char_width_factor
        chars_per_line = max(1, int(width / avg_char_width))
        
        # 単語単位で折り返すことを考慮
        if lang in ["en", "de", "fr", "es", "ru"]:
            # スペースで分割して単語ごとに処理
            words = line.split(' ')
            current_line_length = 0
            line_count = 1
            
            for word in words:
                word_length = len(word) + 1  # スペースを含む
                if current_line_length + word_length > chars_per_line:
                    line_count += 1
                    current_line_length = word_length
                else:
                    current_line_length += word_length
            
            total_lines += line_count
        else:
            # CJK言語などは文字単位で折り返し
            line_count = math.ceil(len(line) / chars_per_line)
            total_lines += line_count
    
    # 行の高さを計算
    line_height = font_size * line_height_factor
    
    # 最終的な高さを返す
    return total_lines * line_height

def adjust_text_box_size(shape, text: str, style: Dict, lang: str = "en", source_lang: str = "ja") -> None:
    """テキストボックスのサイズと位置を調整（改良版）"""
    # 言語ペアごとの幅調整係数
    # 第一キー: ソース言語、第二キー: ターゲット言語
    lang_pair_factors = {
        "ja": {  # 日本語からの翻訳
            "en": 1.3,   # 日本語→英語（英語は日本語より30%広いスペースが必要）
            "zh": 0.9,   # 日本語→中国語（中国語は日本語より10%狭いスペースで済む）
            "ko": 1.0,   # 日本語→韓国語（ほぼ同じ）
            "ru": 1.4,   # 日本語→ロシア語（ロシア語は40%広いスペースが必要）
            "ar": 1.2,   # 日本語→アラビア語
            "de": 1.5,   # 日本語→ドイツ語（ドイツ語は長い単語が多い）
            "fr": 1.4,   # 日本語→フランス語
            "es": 1.3,   # 日本語→スペイン語
        },
        "en": {  # 英語からの翻訳
            "ja": 0.8,   # 英語→日本語（日本語は英語より20%狭いスペースで済む）
            "zh": 0.7,   # 英語→中国語
            "ko": 0.8,   # 英語→韓国語
            "ru": 1.1,   # 英語→ロシア語
            "ar": 0.9,   # 英語→アラビア語
            "de": 1.2,   # 英語→ドイツ語
            "fr": 1.1,   # 英語→フランス語
            "es": 1.1,   # 英語→スペイン語
        },
        # 他の言語ペアも必要に応じて追加
    }
    
    # デフォルトのフォントサイズと最小/最大サイズ
    font_size = 12
    min_font_size = 8  # 最小フォントサイズ
    max_font_size = 24  # 最大フォントサイズ
    
    if "fontSize" in style:
        font_size = style["fontSize"]
        # 最小/最大フォントサイズの制約を設定
        min_font_size = max(8, font_size * 0.7)  # 元のサイズの70%を下限とする
        max_font_size = min(36, font_size * 1.3)  # 元のサイズの130%を上限とする
    
    # 言語ペアに基づく調整係数を取得
    lang_factor = 1.0  # デフォルト値
    if source_lang in lang_pair_factors and lang in lang_pair_factors[source_lang]:
        lang_factor = lang_pair_factors[source_lang][lang]
    
    logger.info(f"言語調整係数: {source_lang}→{lang} = {lang_factor}")
    
    # テキストボックスの元のサイズを保存
    original_width = shape.width
    original_height = shape.height
    
    # テキストの長さに基づいて自動調整を行うかどうかを決定
    text_length_ratio = len(text) / max(1, len(shape.text_frame.text))
    
    # テキスト長の比率に基づいて調整が必要かどうかを判断
    needs_adjustment = text_length_ratio > 1.2 or text_length_ratio < 0.8
    
    # Auto-Fitの設定
    # 0: MSO_AUTO_SIZE.NONE - 自動調整なし
    # 1: MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT - テキストに合わせてシェイプを調整
    # 2: MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE - シェイプに合わせてテキストを調整
    
    # テキスト量が大幅に変わる場合は、テキストボックスのサイズを調整
    if needs_adjustment:
        # テキスト量が多い場合は、ボックスを拡大
        if text_length_ratio > 1.2:
            # 言語係数とテキスト長比率を考慮した調整
            width_adjustment = min(1.5, lang_factor * math.sqrt(text_length_ratio))
            height_adjustment = min(1.5, math.sqrt(text_length_ratio))
            
            # 幅と高さを調整（最大50%増加まで）
            new_width = pixels_to_emu(emu_to_pixels(original_width) * width_adjustment)
            new_height = pixels_to_emu(emu_to_pixels(original_height) * height_adjustment)
            
            # 調整後のサイズを設定
            shape.width = new_width
            shape.height = new_height
            
            # テキストをボックスに合わせる
            shape.text_frame.auto_size = 2  # TEXT_TO_FIT_SHAPE
            
            logger.info(f"テキストボックス拡大: 幅 {width_adjustment:.2f}倍, 高さ {height_adjustment:.2f}倍")
        
        # テキスト量が少ない場合は、フォントサイズを調整
        elif text_length_ratio < 0.8:
            # テキストボックスのサイズはそのままで、フォントサイズを調整
            shape.text_frame.auto_size = 0  # NONE
            
            # フォントサイズを調整（最大30%増加まで）
            font_adjustment = min(1.3, 1.0 / text_length_ratio)
            new_font_size = min(max_font_size, font_size * font_adjustment)
            
            # すべてのテキスト実行に新しいフォントサイズを適用
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(new_font_size)
            
            logger.info(f"フォントサイズ調整: {font_size} → {new_font_size:.1f} (調整係数: {font_adjustment:.2f})")
    else:
        # テキスト量の変化が小さい場合は、標準的な自動調整を使用
        shape.text_frame.auto_size = 1  # SHAPE_TO_FIT_TEXT
        
        # 言語による幅の微調整
        current_width = emu_to_pixels(shape.width)
        adjusted_width = current_width * lang_factor
        shape.width = pixels_to_emu(adjusted_width)
        
        logger.info(f"標準的な自動調整: 幅を{lang_factor:.2f}倍に調整")
    
    # テキストの推定高さを計算し、必要に応じて高さを調整
    estimated_height = estimate_text_height(text, font_size, emu_to_pixels(shape.width), lang)
    current_height = emu_to_pixels(shape.height)
    
    # 推定高さが現在の高さを超える場合、高さを調整
    if estimated_height > current_height * 1.1:  # 10%以上大きい場合のみ調整
        height_ratio = min(1.5, estimated_height / current_height)  # 最大50%まで
        new_height = pixels_to_emu(current_height * height_ratio)
        shape.height = new_height
        logger.info(f"高さ調整: {height_ratio:.2f}倍 ({current_height:.1f} → {current_height * height_ratio:.1f})")
    
    return shape

def adjust_text_frame_properties(text_frame, style: Dict, lang: str = "en") -> None:
    """テキストフレームのプロパティを調整（改良版）"""
    try:
        # 言語ごとの余白調整係数
        margin_factors = {
            "en": {"left": 1.0, "right": 1.0, "top": 1.0, "bottom": 1.0},  # 英語（基準）
            "ja": {"left": 0.8, "right": 0.8, "top": 1.2, "bottom": 1.2},  # 日本語
            "zh": {"left": 0.8, "right": 0.8, "top": 1.2, "bottom": 1.2},  # 中国語
            "ko": {"left": 0.8, "right": 0.8, "top": 1.1, "bottom": 1.1},  # 韓国語
            "ru": {"left": 1.0, "right": 1.0, "top": 1.1, "bottom": 1.1},  # ロシア語
            "ar": {"left": 1.2, "right": 1.2, "top": 1.1, "bottom": 1.1},  # アラビア語
            "de": {"left": 1.0, "right": 1.0, "top": 1.0, "bottom": 1.0},  # ドイツ語
            "fr": {"left": 1.0, "right": 1.0, "top": 1.0, "bottom": 1.0},  # フランス語
            "es": {"left": 1.0, "right": 1.0, "top": 1.0, "bottom": 1.0},  # スペイン語
        }
        
        # 言語に応じた余白係数を取得
        factors = margin_factors.get(lang, {"left": 1.0, "right": 1.0, "top": 1.0, "bottom": 1.0})
        
        # 基本余白値（EMU単位）
        base_margin_h = 73152  # 約0.1インチ（水平方向）
        base_margin_v = 36576  # 約0.05インチ（垂直方向）
        
        # 余白の調整（言語係数を適用）
        text_frame.margin_left = Emu(int(base_margin_h * factors["left"]))
        text_frame.margin_right = Emu(int(base_margin_h * factors["right"]))
        text_frame.margin_top = Emu(int(base_margin_v * factors["top"]))
        text_frame.margin_bottom = Emu(int(base_margin_v * factors["bottom"]))
        
        # スタイル情報から縦方向の配置を取得
        vertical_align = style.get("verticalAlign", "middle").lower()
        
        # 縦方向の配置を設定
        if vertical_align == "top":
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif vertical_align == "bottom":
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:  # デフォルトは中央
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # ワードラップを有効化（スタイルで無効化されていない限り）
        text_frame.word_wrap = style.get("wordWrap", True)
        
        # 行間設定（行間係数）
        line_spacing_factors = {
            "en": 1.0,  # 英語（基準）
            "ja": 1.2,  # 日本語
            "zh": 1.2,  # 中国語
            "ko": 1.15, # 韓国語
            "ru": 1.05, # ロシア語
            "ar": 1.1,  # アラビア語
            "de": 1.0,  # ドイツ語
            "fr": 1.0,  # フランス語
            "es": 1.0,  # スペイン語
        }
        
        # 言語に応じた行間係数を取得
        line_spacing_factor = line_spacing_factors.get(lang, 1.0)
        
        # 段落ごとに行間を設定
        for paragraph in text_frame.paragraphs:
            # 行間を設定（1.0が標準、言語に応じて調整）
            try:
                paragraph.line_spacing = line_spacing_factor
            except Exception as line_spacing_error:
                logger.warning(f"行間設定エラー: {line_spacing_error}")
        
        # 方向が右から左（RTL）の言語の場合
        rtl_languages = ["ar", "he", "fa", "ur"]
        if lang in rtl_languages:
            set_rtl_text_direction(text_frame)
            
        logger.info(f"テキストフレーム調整完了: 言語={lang}, 余白係数={factors}, 行間={line_spacing_factor:.2f}")
    except Exception as e:
        logger.warning(f"テキストフレーム調整エラー: {e}")

def set_rtl_text_direction(text_frame) -> None:
    """テキストフレームを右から左（RTL）に設定"""
    try:
        # xmlタグを直接設定
        p = text_frame._element.get_or_add_p()
        pPr = p.get_or_add_pPr()
        
        # rtl属性を設定
        rtl = OxmlElement('a:rtl')
        rtl.set('val', '1')
        pPr.append(rtl)
    except Exception as e:
        logger.warning(f"RTL設定エラー: {e}")

def create_translated_pptx(original_pptx_path: str, translations: List[Dict], output_path: str):
    """翻訳済みのPPTXファイルを生成"""
    try:
        logger.info(f"翻訳PPTXの生成開始: {original_pptx_path} -> {output_path}")
        
        # 元のプレゼンテーションを読み込み
        prs = Presentation(original_pptx_path)
        
        # プレゼンテーションのメタデータを更新
        if hasattr(prs, 'core_properties'):
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                prs.core_properties.title += " (Translated)"
        
        # 翻訳データの検証
        if not isinstance(translations, list):
            logger.error("翻訳データが無効です: リスト形式ではありません")
            raise ValueError("翻訳データは配列である必要があります")
        
        logger.info(f"スライド数: 元={len(prs.slides)}, 翻訳={len(translations)}")
        
        # 各スライドの翻訳を適用
        for slide_idx, slide_data in enumerate(translations):
            if slide_idx >= len(prs.slides):
                logger.warning(f"スライド {slide_idx + 1} が元のプレゼンテーションに存在しません")
                continue
                
            slide = prs.slides[slide_idx]
            texts = slide_data.get("texts", [])
            target_language = slide_data.get("targetLanguage", "en")
            
            logger.info(f"スライド {slide_idx + 1} の処理 - テキスト要素数: {len(texts)}")
            
            # 既存のテキストボックスを更新
            shape_index = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape_index < len(texts):
                    text_data = texts[shape_index]
                    
                    # 翻訳テキストの取得
                    translation = text_data.get("translation", "") or text_data.get("text", "")
                    
                    # シェイプの種類を確認（テキストボックス、プレースホルダー等）
                    shape_type = None
                    if hasattr(shape, 'shape_type'):
                        shape_type = shape.shape_type
                    
                    logger.info(f"テキスト要素 {shape_index + 1}: タイプ={shape_type}, テキスト長={len(translation)}")
                    
                    # テキストボックスの位置とサイズを更新
                    position = text_data.get("position", {})
                    if position:
                        # スケール係数を元に戻して位置を調整
                        shape.left = pixels_to_emu(position.get("x", 0) * 2)  # スケール係数を元に戻す
                        shape.top = pixels_to_emu(position.get("y", 0) * 2)    # スケール係数を元に戻す
                        
                        # テキストボックスの幅と高さを言語に応じて調整
                        width = pixels_to_emu(position.get("width", 0) * 2)
                        height = pixels_to_emu(position.get("height", 0) * 2)
                        
                        # テキストボックスのサイズを設定
                        shape.width = width
                        shape.height = height
                    
                    # スタイル情報を取得
                    style = text_data.get("style", {})
                    
                    # テキストを更新
                    if shape.text_frame:
                        try:
                            # テキストフレームの基本プロパティを調整
                            adjust_text_frame_properties(shape.text_frame, style, target_language)
                            
                            # 既存のテキストをクリア
                            shape.text_frame.clear()
                            
                            # 翻訳テキストを追加
                            p = shape.text_frame.paragraphs[0]
                            
                            # 段落の配置を設定
                            alignment = text_data.get("alignment", "left")
                            p.alignment = get_alignment(alignment)
                            
                            # テキストとスタイルを追加
                            run = p.add_run()
                            run.text = translation
                            
                            # スタイルを適用
                            if style:
                                apply_text_style(run, style)
                            
                                    # テキストボックスのサイズを調整（ソース言語とターゲット言語を指定）
                            source_lang = slide_data.get("sourceLanguage", "ja")  # デフォルトは日本語
                            adjust_text_box_size(shape, translation, style, target_language, source_lang)
                            
                            logger.info(f"テキスト要素 {shape_index + 1} が更新されました")
                        except Exception as e:
                            logger.error(f"テキスト更新エラー (スライド {slide_idx + 1}, テキスト {shape_index + 1}): {e}")
                    
                    shape_index += 1
        
        # 保存
        prs.save(output_path)
        logger.info(f"翻訳PPTXが生成されました: {output_path}")
        
        return {
            "status": "success",
            "message": "PPTX file generated successfully",
            "path": output_path
        }
        
    except Exception as e:
        logger.error(f"PPTX生成エラー: {e}", exc_info=True)
        return {
            "status": "error",
            "message": f"Failed to generate PPTX: {str(e)}",
            "error": str(e)
        }

if __name__ == "__main__":
    # コマンドライン引数からパスを取得
    if len(sys.argv) != 4:
        print("Usage: python pptx_generator.py <original_pptx_path> <translations_json_path> <output_path>")
        sys.exit(1)
    
    original_pptx_path = sys.argv[1]
    translations_json_path = sys.argv[2]
    output_path = sys.argv[3]
    
    # JSONファイルから翻訳データを読み込む
    try:
        with open(translations_json_path, 'r', encoding='utf-8') as f:
            translations = json.load(f)
    except Exception as e:
        print(f"Error loading translations JSON: {e}")
        sys.exit(1)
    
    # PPTXを生成
    result = create_translated_pptx(original_pptx_path, translations, output_path)
    
    # 結果を出力
    print(json.dumps(result))