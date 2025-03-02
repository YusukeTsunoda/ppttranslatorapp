#!/usr/bin/env python3
import sys
import json
import os
from pptx import Presentation
from typing import List, Dict, Any, Tuple
import subprocess
from pdf2image import convert_from_path
import logging

logger = logging.getLogger(__name__)

class PPTXParser:
    @staticmethod
    def convert_to_png(pptx_path: str, output_dir: str) -> Tuple[List[str], Tuple[int, int]]:
        """
        PPTXの各スライドをPNG画像に変換する
        
        Args:
            pptx_path (str): PPTXファイルのパス
            output_dir (str): 出力ディレクトリ
        
        Returns:
            tuple[List[str], tuple[int, int]]: 画像パスのリストと画像サイズ
        """
        try:
            # 出力ディレクトリが存在しない場合は作成
            os.makedirs(output_dir, exist_ok=True)
            
            # PPTXをPDFに変換
            pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(output_dir, f"{pptx_name}.pdf")
            
            logger.info(f"Converting PPTX to PDF: {pptx_path} -> {pdf_path}")
            
            if sys.platform == "darwin":  # macOS
                result = subprocess.run(
                    ["soffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
                    capture_output=True,
                    text=True
                )
                if result.returncode != 0:
                    logger.error(f"Error converting to PDF: {result.stderr}")
                    return [], (0, 0)
            else:  # Linux
                result = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
                    capture_output=True,
                    text=True
                )
                if result.returncode != 0:
                    logger.error(f"Error converting to PDF: {result.stderr}")
                    return [], (0, 0)

            if not os.path.exists(pdf_path):
                logger.error(f"PDF file not found at: {pdf_path}")
                return [], (0, 0)
                
            logger.info(f"Converting PDF to images: {pdf_path}")
                
            # PDFを画像に変換（16:9のアスペクト比を保持）
            try:
                # 720x405のサイズ（16:9）で変換
                target_width = 720
                target_height = int(target_width * 9 / 16)  # 405px
                images = convert_from_path(pdf_path, size=(target_width, target_height))
                
                if not images:
                    return [], (0, 0)
                    
                # 実際の画像サイズを取得
                actual_width = images[0].width
                actual_height = images[0].height
                
                logger.info(f"Generated image size: {actual_width}x{actual_height}")
                
            except Exception as e:
                logger.error(f"Error converting PDF to images: {str(e)}")
                return [], (0, 0)
            
            image_paths = []
            for i, image in enumerate(images):
                image_path = os.path.join(output_dir, f"slide_{i+1}.png")
                image.save(image_path, "PNG")
                rel_path = f"slide_{i+1}.png"
                image_paths.append(rel_path)
                logger.info(f"Saved image: {image_path}")
            
            try:
                os.remove(pdf_path)
                logger.info(f"Removed temporary PDF: {pdf_path}")
            except Exception as e:
                logger.warning(f"Could not delete temporary PDF: {str(e)}")
            
            return image_paths, (actual_width, actual_height)
            
        except Exception as e:
            logger.error(f"Error converting to PNG: {str(e)}")
            return [], (0, 0)

    @staticmethod
    def convert_coordinates(x: float, y: float, width: float, height: float, 
                          slide_width: float, slide_height: float,
                          image_width: float, image_height: float) -> dict:
        """
        PowerPointの座標系（EMU）をピクセル座標に変換
        
        Args:
            x, y: 元の座標（EMU）
            width, height: 元のサイズ（EMU）
            slide_width: スライドの幅（EMU）
            slide_height: スライドの高さ（EMU）
            image_width: 画像の幅（ピクセル）
            image_height: 画像の高さ（ピクセル）
        
        Returns:
            dict: 変換後の座標とサイズ
        """
        # スライドと画像のアスペクト比を計算
        slide_aspect = slide_width / slide_height
        image_aspect = image_width / image_height
        
        # 基本のスケーリング係数を計算（幅に基づく）
        scale = image_width / slide_width
        
        # 座標を変換
        scaled_x = x * scale
        scaled_y = y * scale
        scaled_width = width * scale
        scaled_height = height * scale
        
        # アスペクト比の違いによる調整
        if slide_aspect != image_aspect:
            # 高さの差を計算
            expected_height = slide_height * scale
            height_diff = image_height - expected_height
            # Y座標を調整（中央揃え）
            scaled_y += height_diff / 2
        
        return {
            "x": round(scaled_x),
            "y": round(scaled_y),
            "width": round(scaled_width),
            "height": round(scaled_height)
        }

    @staticmethod
    def extract_text_from_shape(shape):
        if hasattr(shape, "text"):
            return shape.text.strip()
        return ""

    @staticmethod
    def parse_pptx(file_path: str, output_dir: str) -> List[Dict[str, Any]]:
        try:
            logger.info(f"Processing PPTX file: {file_path}")
            prs = Presentation(file_path)
            slides = []
            
            # スライドを画像に変換
            image_paths, (image_width, image_height) = PPTXParser.convert_to_png(file_path, output_dir)
            
            if not image_paths:
                logger.error("No images were generated")
                return []
                
            for slide_index, slide in enumerate(prs.slides):
                texts = []
                
                # スライドのサイズを取得
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                for shape in slide.shapes:
                    text = PPTXParser.extract_text_from_shape(shape)
                    if text:
                        # 座標を変換
                        position = PPTXParser.convert_coordinates(
                            shape.left, shape.top, 
                            shape.width, shape.height,
                            slide_width, slide_height,
                            image_width, image_height
                        )
                        texts.append({
                            "text": text,
                            "type": "text",
                            "position": position
                        })
                
                slides.append({
                    "index": slide_index,
                    "texts": texts,
                    "image_path": image_paths[slide_index] if slide_index < len(image_paths) else None
                })
            
            return slides
            
        except Exception as e:
            logger.error(f"Error parsing PPTX: {str(e)}")
            return []

    @staticmethod
    def update_pptx_with_translations(
        original_file: str,
        output_file: str,
        translations: Dict[int, List[Dict[str, str]]]
    ) -> bool:
        """
        翻訳されたテキストでPPTXファイルを更新する
        
        Args:
            original_file (str): 元のPPTXファイルパス
            output_file (str): 出力するPPTXファイルパス
            translations (Dict[int, List[Dict[str, str]]]): 
                スライドインデックスごとの翻訳テキスト情報
        
        Returns:
            bool: 成功した場合True
        """
        try:
            prs = Presentation(original_file)
            
            for slide_index, slide in enumerate(prs.slides):
                if slide_index not in translations:
                    continue
                    
                slide_translations = translations[slide_index]
                shape_map = {shape.shape_id: shape for shape in slide.shapes 
                            if hasattr(shape, "text")}
                
                for trans in slide_translations:
                    shape_id = trans.get("shape_id")
                    if shape_id in shape_map:
                        shape = shape_map[shape_id]
                        shape.text = trans["translated_text"]
            
            prs.save(output_file)
            return True
            
        except Exception as e:
            logger.error(f"Error updating PPTX: {str(e)}")
            return False 