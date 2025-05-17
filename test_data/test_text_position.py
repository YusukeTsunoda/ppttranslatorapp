#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
テキスト位置調整機能のテストスクリプト
異なる言語ペアでのテキストボックスサイズ調整をテストします
"""

import sys
import os
import json
import logging
from pptx import Presentation
import argparse
from pathlib import Path

# 親ディレクトリをパスに追加
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from lib.python.pptx_generator import create_translated_pptx, adjust_text_box_size

# ロギングの設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('text_position_test')

def generate_test_translations(original_pptx_path, language_pairs):
    """
    テスト用の翻訳データを生成します
    
    Args:
        original_pptx_path: 元のPPTXファイルのパス
        language_pairs: テストする言語ペアのリスト [(source_lang, target_lang), ...]
    
    Returns:
        翻訳データのリスト
    """
    # 元のPPTXを読み込み
    prs = Presentation(original_pptx_path)
    
    # 各言語ペアごとの翻訳データを生成
    all_translations = {}
    
    for source_lang, target_lang in language_pairs:
        translations = []
        
        # 各スライドの翻訳データを生成
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {"slideIndex": slide_idx, "texts": []}
            
            # テキストボックスの内容を取得して翻訳データを生成
            shape_index = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    original_text = shape.text_frame.text
                    
                    # テスト用の翻訳テキストを生成（言語ペアに応じて特性を模倣）
                    translated_text = generate_mock_translation(original_text, source_lang, target_lang)
                    
                    # テキストデータを追加
                    text_data = {
                        "index": shape_index,
                        "text": original_text,
                        "translation": translated_text,
                        "position": {
                            "x": shape.left / 9525 / 2,  # EMUからピクセルに変換して調整
                            "y": shape.top / 9525 / 2,
                            "width": shape.width / 9525 / 2,
                            "height": shape.height / 9525 / 2
                        },
                        "style": {
                            "fontSize": 12,  # デフォルト値
                            "fontFamily": "Arial",
                            "alignment": "left"
                        }
                    }
                    
                    slide_data["texts"].append(text_data)
                    shape_index += 1
            
            translations.append(slide_data)
        
        # 言語ペアごとの翻訳データを保存
        pair_key = f"{source_lang}_to_{target_lang}"
        all_translations[pair_key] = translations
    
    return all_translations

def generate_mock_translation(text, source_lang, target_lang):
    """
    テスト用の模擬翻訳を生成します
    言語ペアの特性に基づいて、テキスト長を調整します
    
    Args:
        text: 元のテキスト
        source_lang: ソース言語
        target_lang: ターゲット言語
    
    Returns:
        模擬翻訳テキスト
    """
    # 言語ペアごとのテキスト長変化率
    length_factors = {
        "ja_to_en": 1.3,   # 日本語→英語（英語は日本語より30%長い）
        "en_to_ja": 0.8,   # 英語→日本語（日本語は英語より20%短い）
        "ja_to_zh": 0.9,   # 日本語→中国語
        "en_to_zh": 0.7,   # 英語→中国語
        "ja_to_de": 1.5,   # 日本語→ドイツ語（ドイツ語は長い単語が多い）
        "en_to_de": 1.2,   # 英語→ドイツ語
        "ja_to_fr": 1.4,   # 日本語→フランス語
        "en_to_fr": 1.1,   # 英語→フランス語
        "ja_to_ru": 1.4,   # 日本語→ロシア語
        "en_to_ru": 1.1,   # 英語→ロシア語
    }
    
    # 言語ペアのキー
    pair_key = f"{source_lang}_to_{target_lang}"
    
    # 言語ペアの長さ係数を取得（デフォルトは1.0）
    length_factor = length_factors.get(pair_key, 1.0)
    
    # 言語ペアごとの特性を模倣した翻訳テキストを生成
    if target_lang == "en":
        # 英語の場合は、単語の長さを調整
        words = ["This", "is", "a", "test", "translation", "for", "English", "text", 
                "with", "some", "longer", "words", "to", "simulate", "real", "content"]
        word_count = max(3, int(len(text) * length_factor / 5))  # 平均5文字/単語と仮定
        mock_text = " ".join(words[:word_count % len(words)] + words * (word_count // len(words)))
        return mock_text[:int(len(text) * length_factor)]
    
    elif target_lang == "de":
        # ドイツ語の場合は、長い複合語を含める
        words = ["Das", "ist", "ein", "Test", "für", "Deutsch", "mit", "einigen", 
                "Kompositumwörtern", "Zeitungsredakteur", "Donaudampfschifffahrtsgesellschaft", 
                "Rindfleischetikettierungsüberwachungsaufgabenübertragungsgesetz"]
        word_count = max(2, int(len(text) * length_factor / 8))  # 平均8文字/単語と仮定
        mock_text = " ".join(words[:word_count % len(words)] + words * (word_count // len(words)))
        return mock_text[:int(len(text) * length_factor)]
    
    elif target_lang == "ja":
        # 日本語の場合は、漢字とひらがなを混ぜる
        chars = "これはテスト用の日本語テキストです。漢字とひらがなが混ざっています。長さを調整して実際の翻訳を模倣します。"
        return chars[:int(len(text) * length_factor)]
    
    elif target_lang == "zh":
        # 中国語の場合は、漢字のみ
        chars = "这是一个测试用的中文文本。它包含了汉字。我们调整长度来模拟实际的翻译。中文通常比英文短。"
        return chars[:int(len(text) * length_factor)]
    
    elif target_lang == "ru":
        # ロシア語の場合は、キリル文字を使用
        chars = "Это тестовый текст на русском языке. Он содержит кириллические символы. Мы регулируем длину для имитации реального перевода."
        return chars[:int(len(text) * length_factor)]
    
    elif target_lang == "fr":
        # フランス語の場合は、アクセント記号を含む
        words = ["C'est", "un", "texte", "de", "test", "en", "français", "avec", "des", "accents", 
                "comme", "é", "è", "ê", "à", "ù", "ç", "pour", "simuler", "une", "traduction", "réelle"]
        word_count = max(3, int(len(text) * length_factor / 5))
        mock_text = " ".join(words[:word_count % len(words)] + words * (word_count // len(words)))
        return mock_text[:int(len(text) * length_factor)]
    
    else:
        # その他の言語の場合は、単純に長さを調整
        return "X" * int(len(text) * length_factor)

def run_test(original_pptx_path, output_dir, language_pairs):
    """
    テストを実行します
    
    Args:
        original_pptx_path: 元のPPTXファイルのパス
        output_dir: 出力ディレクトリ
        language_pairs: テストする言語ペアのリスト [(source_lang, target_lang), ...]
    """
    # 出力ディレクトリを作成
    os.makedirs(output_dir, exist_ok=True)
    
    # テスト用の翻訳データを生成
    all_translations = generate_test_translations(original_pptx_path, language_pairs)
    
    # 各言語ペアでテスト
    for pair_key, translations in all_translations.items():
        source_lang, target_lang = pair_key.split("_to_")
        
        # 翻訳データをJSONファイルに保存
        translations_json_path = os.path.join(output_dir, f"{pair_key}_translations.json")
        with open(translations_json_path, 'w', encoding='utf-8') as f:
            json.dump(translations, f, ensure_ascii=False, indent=2)
        
        # 翻訳PPTXを生成
        output_pptx_path = os.path.join(output_dir, f"{pair_key}_translated.pptx")
        
        # 翻訳を実行
        logger.info(f"言語ペア {pair_key} のテストを実行中...")
        result = create_translated_pptx(original_pptx_path, translations, output_pptx_path)
        
        if result.get("status") == "success":
            logger.info(f"言語ペア {pair_key} のテスト成功: {output_pptx_path}")
        else:
            logger.error(f"言語ペア {pair_key} のテスト失敗: {result.get('message')}")
    
    logger.info("すべてのテストが完了しました。")
    return all_translations

def main():
    parser = argparse.ArgumentParser(description='テキスト位置調整機能のテスト')
    parser.add_argument('--input', '-i', required=True, help='テスト用の元PPTXファイルのパス')
    parser.add_argument('--output-dir', '-o', default='./test_results', help='テスト結果の出力ディレクトリ')
    args = parser.parse_args()
    
    # テストする言語ペア
    language_pairs = [
        ("ja", "en"),  # 日本語→英語
        ("en", "ja"),  # 英語→日本語
        ("ja", "zh"),  # 日本語→中国語
        ("en", "zh"),  # 英語→中国語
        ("ja", "de"),  # 日本語→ドイツ語
        ("en", "de"),  # 英語→ドイツ語
        ("ja", "fr"),  # 日本語→フランス語
        ("en", "fr"),  # 英語→フランス語
        ("ja", "ru"),  # 日本語→ロシア語
        ("en", "ru"),  # 英語→ロシア語
    ]
    
    # テストを実行
    run_test(args.input, args.output_dir, language_pairs)

if __name__ == "__main__":
    main()
